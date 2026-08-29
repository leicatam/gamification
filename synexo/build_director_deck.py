#!/usr/bin/env python3
"""Executive presentation (30 min) -> GenApep_Executive_Presentation.pptx.

For the directors of WBI and Eyesel: the IPO project, First Vital, routes,
recommended shareholding structure, WBI-shareholder lens, decisions.
16 slides at ~2 min each. House deck style (matches build_pptx.py).
Figures cited from the JV Financial Plan (rev 29 Aug 2026), the GLP-1
proforma, the Stage 2 Integration Framework and the DPW licensing framework.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

DIR = "/home/user/gamification/synexo/"
NAVY = RGBColor(0x1B, 0x33, 0x55); ACCENT = RGBColor(0x2E, 0x86, 0xC1)
TEAL = RGBColor(0x1A, 0xA1, 0x9C); GOLD = RGBColor(0xE0, 0x9E, 0x2B)
GREY = RGBColor(0x5A, 0x5A, 0x5A); LIGHT = RGBColor(0xEC, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); RED = RGBColor(0xC0, 0x39, 0x2B)
SKY = RGBColor(0xCF, 0xDD, 0xEA); DARK2 = RGBColor(0x24, 0x44, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
IN = Inches
CURRENT = [0]


def slide():
    CURRENT[0] += 1
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0, shadow=False):
    sp = s.shapes.add_shape(shape, IN(l), IN(t), IN(w), IN(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = shadow
    return sp


def fill_text(sp, lines, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"


def textbox(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
        p.space_after = Pt(4)
    return tb


def header(s, title, kicker=None):
    rect(s, 0, 0, 13.333, 1.2, NAVY)
    rect(s, 0.55, 0.34, 0.14, 0.52, ACCENT)
    textbox(s, 0.85, 0.18, 11.6, 0.9, [(title, 27, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        textbox(s, 0.87, 0.86, 11.6, 0.3, [(kicker, 14, False, SKY)])
    textbox(s, 12.4, 7.04, 0.8, 0.35, [(str(CURRENT[0]), 12, False, GREY)], align=PP_ALIGN.RIGHT)


def bullets(s, items, l=0.9, t=1.6, w=11.6, h=5.2, size=18):
    tb = s.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("●  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size - lvl * 3); r.font.name = "Calibri"
        r.font.color.rgb = NAVY if lvl == 0 else GREY; r.font.bold = False
        p.space_after = Pt(11 if lvl == 0 else 5)
    return tb


def card(s, l, t, w, h, title, body, fill=WHITE, tcol=NAVY, bcol=GREY, tsz=17, bsz=14,
         line=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = rect(s, l, t, w, h, fill, line=line, line_w=1.5, shape=shape, shadow=True)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(10); tf.margin_right = Pt(10); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title; r.font.size = Pt(tsz); r.font.bold = True
    r.font.color.rgb = tcol; r.font.name = "Calibri"
    if body:
        p2 = tf.add_paragraph(); p2.space_before = Pt(4); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(bsz); r2.font.color.rgb = bcol
        r2.font.name = "Calibri"
    return sp


def stat(s, l, t, w, h, big, label, fill=NAVY, bigcol=WHITE, labcol=SKY, bigsz=34, labsz=13):
    sp = rect(s, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    fill_text(sp, [(big, bigsz, True, bigcol), (label, labsz, False, labcol)])
    return sp


def table(s, headers, rows, l, t, w, col_w, hdr_fs=13, fs=12.5, hl_rows=(), row_h=None):
    gt = s.shapes.add_table(len(rows) + 1, len(headers), IN(l), IN(t), IN(w), IN(0.4 * (len(rows) + 1))).table
    for i, cw in enumerate(col_w):
        gt.columns[i].width = IN(cw)
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            for r in p.runs: r.font.size = Pt(hdr_fs); r.font.bold = True; r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = gt.cell(i, j); c.text = str(v)
            c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = NAVY
                    r.font.bold = (j == 0 or (i - 1) in hl_rows)
    return gt


def chart_fonts(chart, size=13):
    try:
        chart.value_axis.tick_labels.font.size = Pt(size)
        chart.category_axis.tick_labels.font.size = Pt(size)
    except Exception:
        pass
    if chart.has_legend:
        chart.legend.font.size = Pt(size)


# figures
POOL = 0.075; EY_POST = 0.51 * (1 - POOL); WBI_POST = 0.49 * (1 - POOL)
GA_MID = 35.0
def block(fvh_v): return GA_MID / (GA_MID + fvh_v)

# ============================================================ 1 TITLE
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 13.333, 0.28, ACCENT); rect(s, 0, 7.22, 13.333, 0.28, ACCENT)
for (cx, cy, d, col) in [(11.4, 1.5, 2.6, DARK2), (12.2, 5.6, 1.8, DARK2), (10.6, 6.4, 1.0, ACCENT)]:
    rect(s, cx, cy, d, d, col, shape=MSO_SHAPE.OVAL)
rect(s, 0.95, 2.35, 0.2, 1.9, ACCENT)
textbox(s, 1.35, 2.25, 10.5, 1.3, [("GENAPEP — THE ROAD TO A US LISTING", 44, True, WHITE)])
textbox(s, 1.4, 3.65, 10.5, 0.9, [("Director briefing: the Joint Venture, First Vital, and the choice of route",
                                   24, False, SKY)])
textbox(s, 1.4, 4.55, 10.5, 0.6, [("Executive presentation · ~30 minutes · read before the document pack",
                                   16, False, RGBColor(0x9F, 0xB6, 0xCB))])
textbox(s, 1.4, 6.35, 11.5, 0.7, [("PRIVATE & CONFIDENTIAL — for the directors of WBI and Eyesel only. "
                                   "Not for release to First Vital. Illustrative planning figures; not advice.",
                                   13, True, GOLD)])

# ============================================================ 2 WHY WE'RE HERE
s = slide(); header(s, "Why we are here", "Thirty minutes before you read eight documents")
bullets(s, [
    "WBI and Eyesel signed an MOU on 12 August 2026 to combine into one group — GenApep — and take it to a US public market.",
    "Two new elements have arrived since: First Vital (a US company with a stated SEC filing record and a GLP-1 care programme) proposes to be the listing vehicle, and DPW (under the listed Richards Group) is reviewing a licence to GenApep's technology.",
    "The document pack runs to eight instruments and two workbooks. This presentation gives you the structure, the numbers, the pros and cons, and the decisions — so the pack confirms rather than surprises.",
    ("The ask today: agree the implementation order, note the risks, and approve the next actions on the final slide.", 0),
])
stat(s, 9.0, 5.6, 3.65, 1.3, "8 + 2", "documents + workbooks distilled here", fill=ACCENT, bigsz=30)

# ============================================================ 3 ONE PICTURE
s = slide(); header(s, "The project in one picture", "Two stages; the second is a choice, not a commitment")
card(s, 0.7, 1.7, 5.9, 2.5, "STAGE 1 — form the Joint Venture (now)",
     "GenApep Holdings (Cayman) → GenApep Korea → acquires 100% of WBI and 100% of Eyesel.\n"
     "Eyesel shareholders 51% · WBI shareholders 49% · tech sweat pool 5–10% dilutes both pro rata.\n"
     "Governed by the 12-Aug MOU + Amendment No. 1.", line=ACCENT, tsz=18, bsz=14)
card(s, 6.85, 1.7, 5.8, 2.5, "STAGE 2 — choose the route to market (later, gated)",
     "Route A (default): GenApep's own US offering — SEC confidential draft by end-2027, listing 2028.\n"
     "Route B (option): combine with First Vital, a company with an existing SEC record — only "
     "through the machinery in Addendum A.", line=GOLD, tsz=18, bsz=14)
a = rect(s, 6.35, 2.75, 0.5, 0.45, TEAL, shape=MSO_SHAPE.RIGHT_ARROW)
band = rect(s, 0.7, 4.5, 11.95, 1.0, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
fill_text(band, [("The PCAOB audit of both companies is the critical path under EVERY route — "
                  "so forming the JV first costs no calendar time.", 17, True, NAVY)])
card(s, 0.7, 5.75, 11.95, 1.15, "Why the order matters",
     "Settle the 51:49 economics between ourselves first, then negotiate with First Vital as ONE "
     "party holding audited accounts — instead of two unaudited companies and two individuals "
     "inside someone else's cap table.", line=RED, tsz=16, bsz=14)

# ============================================================ 4 STAGE 1 ESSENTIALS
s = slide(); header(s, "Stage 1 — the Joint Venture, in essentials", "12-Aug MOU as amended by Amendment No. 1")
card(s, 0.7, 1.6, 3.85, 2.6, "Structure",
     "Cayman holdco → GenApep Korea → 100% WBI + 100% Eyesel. Korean operations undisturbed; "
     "one group for audit and listing.", line=ACCENT)
card(s, 4.75, 1.6, 3.85, 2.6, "Economics",
     "Eyesel 51 : WBI 49, both diluted by a 5–10% technology pool (Keith · Sidney · Charles + team). "
     "Mandatory equalisation if either side's figures move >15%.", line=TEAL)
card(s, 8.8, 1.6, 3.85, 2.6, "Governance",
     "One director each + a Third Director (interests disclosed; deadlock fallback fixes the pool "
     "at 7.5% if directors cannot agree).", line=GOLD)
card(s, 0.7, 4.4, 5.9, 2.3, "Amendment No. 1 — why it must be signed first",
     "Nine fixes needed under every scenario: consideration flows to SHAREHOLDER BODIES pro rata, "
     "never to individuals · equalisation becomes mandatory · technology chain-of-title becomes a "
     "condition precedent · open items get owners and dates.", line=RED, bsz=13.5)
card(s, 6.85, 4.4, 5.8, 2.3, "Addendum A — sign when Route B is on the table",
     "Route selection becomes a board decision · GenApep Korea is the ONLY permitted acquisition "
     "vehicle · any Stage-2 deal must carry every shareholder through pro rata as one block · the "
     "sweat pool survives any change of vehicle.", line=DARK2, bsz=13.5)

# ============================================================ 5 NUMBERS TODAY
s = slide(); header(s, "The numbers today", "JV Financial Plan (rev. 29 Aug 2026) — planning basis, Eyesel figures placeholder")
cd = CategoryChartData(); cd.categories = ["FY2025A", "2026P", "2027P", "2028P", "2029P", "2030P"]
cd.add_series("Revenue", (11.43, 16.39, 25.31, 31.70, 43.82, 58.91))
cd.add_series("Operating profit", (0.76, 2.40, 4.54, 5.95, 8.86, 12.51))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, IN(0.7), IN(1.6), IN(7.6), IN(5.0), cd)
ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.value_axis.has_title = True; ch.value_axis.axis_title.text_frame.text = "US$M"
ch.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
pl = ch.plots[0]; pl.gap_width = 80
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = ACCENT
ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb = NAVY
chart_fonts(ch, 13)
stat(s, 8.6, 1.6, 4.05, 1.25, "$58.9M", "2030 plan revenue (base case)", fill=ACCENT, bigsz=30)
stat(s, 8.6, 3.0, 4.05, 1.25, "$30–40M", "combined pre-JV valuation range", fill=NAVY, bigsz=30)
card(s, 8.6, 4.4, 4.05, 2.2, "What drives it",
     "MM Studio ramp $2.9M → $36.3M is the swing factor (WBI-side). Eyesel base $8.0M placeholder "
     "+10%/yr — updated financials due within 30 days and the 51:49 rests on them.", line=GOLD, bsz=13)

# ============================================================ 6 TWO ROUTES
s = slide(); header(s, "Stage 2 — two routes to being public", "Route A is the default; Route B is an option with gates")
card(s, 0.7, 1.65, 5.9, 3.5, "ROUTE A — our own offering (default)",
     "GenApep Holdings files its own registration.\n"
     "• Confidential SEC draft by end-2027; listing 2028.\n"
     "• Raises primary capital (illustrative $20M at $80M pre-money).\n"
     "• Foreign private issuer: F-1, IFRS, 20-F — as budgeted.\n"
     "• Fully in our control; no counterparty risk.", line=ACCENT, tsz=18, bsz=14.5)
card(s, 6.85, 1.65, 5.8, 3.5, "ROUTE B — combine with First Vital (option)",
     "GenApep enters a company that already reports to the SEC.\n"
     "• Potentially faster to reporting status — but that is NOT an exchange listing.\n"
     "• Raises no money by itself; financing must be planned separately.\n"
     "• Likely domestic issuer: US GAAP, 10-K/10-Q — unbudgeted cost.\n"
     "• Brings the GLP-1 programme and a US clinical channel.", line=GOLD, tsz=18, bsz=14.5)
band = rect(s, 0.7, 5.45, 11.95, 1.3, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
fill_text(band, [("Addendum A keeps BOTH routes open until the boards choose — route selection is a "
                  "Reserved Matter requiring both directors.", 16, True, NAVY),
                 ("No step may foreclose either route without Board approval.", 13, False, GREY)])

# ============================================================ 7 FIRST VITAL
s = slide(); header(s, "First Vital — what it brings", "Stated SEC record + an operating GLP-1 programme (proforma basis)")
card(s, 0.7, 1.6, 5.9, 2.5, "The GLP-1 Bridge programme",
     "CPT-reimbursed care: assessment → care plan → monthly remote monitoring at $168.92/patient/"
     "month. Targeting 2,500 patients by Dec-2026, 10,000 by Dec-2027.", line=TEAL, bsz=14)
stat(s, 0.7, 4.3, 2.85, 1.3, "$2.2M", "proforma revenue H2-2026", fill=TEAL, bigsz=28)
stat(s, 3.75, 4.3, 2.85, 1.3, "$16.7M", "proforma revenue 2027", fill=ACCENT, bigsz=28)
stat(s, 0.7, 5.75, 2.85, 1.3, "$1.59M/mo", "recurring MRR exiting 2027", fill=NAVY, bigsz=24)
stat(s, 3.75, 5.75, 2.85, 1.3, "~$3,788", "revenue per enrolled patient", fill=DARK2, bigsz=26)
card(s, 6.85, 1.6, 5.8, 2.5, "Why it fits GenApep",
     "A GLP-1 patient base is a natural US channel for peptide products: GI & metabolic support, "
     "muscle preservation, GLP-1-associated hair loss — designed by the AI platform, made by "
     "Eyesel under GMP.", line=GOLD, bsz=14)
card(s, 6.85, 4.3, 5.8, 2.75, "Read the numbers with care",
     "• Proforma is GROSS REIMBURSEMENT only — no clinician, billing or collection costs.\n"
     "• Programme began enrolling 1 July 2026 — two months of history.\n"
     "• SEC filing record is stated, not yet verified on EDGAR.\n"
     "• July–August actuals are the cheapest, most informative diligence available.", line=RED, bsz=13.5)

# ============================================================ 8 ROUTE B PROS
s = slide(); header(s, "Route B — the case for", "What a verified First Vital adds")
items = [
    ("Time", "An existing, current filing record cuts the time and execution risk of becoming a reporting company.", ACCENT),
    ("Channel", "10,000 targeted GLP-1 patients = a US clinical route to market GenApep cannot quickly build alone.", TEAL),
    ("Revenue", "A recurring-revenue care programme (MRR $1.59M/month exiting 2027, proforma) alongside our product revenue.", GOLD),
    ("Synergy", "GLP-1 GI support + hair-loss management are exactly the categories GenApep's peptide/exosome stack targets.", DARK2),
    ("Infrastructure", "US corporate platform and an existing shareholder base — float the JV would otherwise create from zero.", NAVY),
]
for i, (t, d, col) in enumerate(items):
    y = 1.65 + i * 1.06
    badge = rect(s, 0.8, y, 1.7, 0.85, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    fill_text(badge, [(t, 16, True, WHITE)])
    textbox(s, 2.75, y + 0.06, 9.8, 0.85, [(d, 15.5, False, NAVY)], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 9 ROUTE B CONS
s = slide(); header(s, "Route B — the risks", "Each has a gate in Addendum A / the diligence programme")
rows = [
    ["Filing status unverified", "No identifiable EDGAR record surfaced by name search; a Reg-A filing would not make FVH a reporting company", "Gate zero: CIK + full EDGAR index, counsel confirmation"],
    ["Two months of history", "GLP-1 enrolment began 1 Jul 2026; the 10,000-patient ramp is a plan", "Jul–Aug actual enrolments, claims, denials — this week"],
    ["Revenue-only proforma", "No cost build; 20% margin is an assumption; reimbursement ≠ cash", "Cost build + collection-rate before any valuation"],
    ["Issuer & accounting", "US parent likely = domestic issuer: US GAAP, 10-K/10-Q, proxy rules", "Written classification opinion before route selection"],
    ["Shell / liabilities", "Possible shell-company resale restrictions; unknown liabilities", "Shell status on EDGAR; corporate & legal diligence"],
    ["No capital raised", "A combination raises nothing; group could emerge public but underfunded", "Separate financing plan before selection"],
    ["Related-party optics", "Principals could sit on both sides of the transaction", "Independent valuations + fairness opinion + disinterested approval"],
]
table(s, ["Risk", "Why it matters", "The gate"], rows, 0.7, 1.6, 11.95, [2.25, 5.4, 4.3], fs=11.5, hdr_fs=12.5)

# ============================================================ 10 THE CONFLICT
s = slide(); header(s, "The conflict to resolve first", "The draft IPO MOU cannot be signed alongside the JV MOU")
rows = [
    ["Who gets the shares", "WBI & Eyesel SHAREHOLDER BODIES, pro rata", "Two named individuals (40% + 35%)"],
    ["The split", "Eyesel 51 : WBI 49, less sweat pool", "≈ 53 : 47 the other way; no pool at all"],
    ["Acquisition vehicle", "GenApep Korea (sole permitted vehicle)", "FVH Korea — a competing vehicle"],
    ["When shares move", "Issuance and transfer simultaneous at closing", "75% issued BEFORE the companies transfer"],
    ["Eyesel principal / revenue", "Mr. Park · US$8.0M (plan placeholder)", "Mr. Kim · ~US$15M"],
]
table(s, ["", "JV MOU + Amendment (agreed)", "Draft IPO MOU (as written)"], rows,
      0.7, 1.6, 11.95, [2.6, 4.7, 4.65], fs=12.5, hdr_fs=13)
band = rect(s, 0.7, 4.95, 11.95, 1.7, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
fill_text(band, [("The fix is sequencing, not abandonment.", 18, True, NAVY),
                 ("Form the Joint Venture first, then approach First Vital as one block with one agreed "
                  "internal split. Same commercial outcome, none of the governance defects — and zero "
                  "calendar cost, because the audit runs on the critical path either way.", 14.5, False, GREY)])

# ============================================================ 11 RECOMMENDED STRUCTURE
s = slide(); header(s, "Recommended shareholding structure", "FVH's structure is undetermined — set it by value, enter as one block")
card(s, 0.7, 1.55, 5.9, 2.05, "The principle",
     "One value-based exchange ratio: independent valuations of GenApep ($30–40M planning range) "
     "and of a VERIFIED First Vital + a written fairness opinion. FVH issues shares to GenApep "
     "Holdings' shareholders pro rata — never to individuals.", line=ACCENT, bsz=13.5)
card(s, 0.7, 3.75, 5.9, 1.85, "Technology credit — the DPW licence",
     "DPW (under listed Richards Group) is reviewing a licence to GenApep technology. An "
     "arm's-length licence evidences the $5–10M AI.pep value now carried at $0 — and adds royalty "
     "income. Gate: chain-of-title first; keep exclusivity narrow.", line=GOLD, bsz=13)
rows = [
    ["FVH verified $5M", "88%", "41.3%", "39.7%", "6.6%", "12%"],
    ["FVH verified $10M", "78%", "36.7%", "35.3%", "5.8%", "22%"],
    ["FVH verified $13M", "73%", "34.4%", "33.0%", "5.5%", "27%"],
]
textbox(s, 6.85, 1.55, 5.8, 0.4, [("Illustrative cap table (GenApep $35M mid · 7.5% pool)", 15, True, NAVY)])
table(s, ["Scenario", "GA block", "Eyesel", "WBI", "Pool", "FVH"], rows,
      6.85, 2.05, 5.8, [1.75, 0.95, 0.85, 0.8, 0.7, 0.75], fs=12, hdr_fs=12, hl_rows=(1,))
card(s, 6.85, 4.1, 5.8, 1.5, "Why this beats 40/35/25",
     "The block percentage is the only number to negotiate; the internal 51:49-less-pool split "
     "then divides automatically. Simpler to agree, to explain, and to support with a fairness "
     "opinion.", line=TEAL, bsz=13)
band = rect(s, 0.7, 5.85, 11.95, 0.85, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
fill_text(band, [("Price test: the draft's percentages imply paying $10–13M for FVH. Would FVH be worth "
                  "that if it were private? Verify before a single share moves.", 14.5, True, NAVY)])

# ============================================================ 12 WBI SHAREHOLDER LENS
s = slide(); header(s, "What this means for WBI shareholders", "Your 49% of GenApep, traced through both routes")
steps = [("Today", "49%", "of GenApep · ≈ $17.2M at $35M mid", ACCENT),
         ("After pool (7.5%)", "45.3%", "≈ $15.9M — before technology credit", TEAL),
         ("Route B (FVH $10M)", "35.3%", "of the LISTED company · $28–35M if market values group at $80–100M", GOLD),
         ("Route A post-IPO", "36.3%", "≈ $36.3M at $80M pre + $20M raise (workbook)", NAVY)]
for i, (t, big, d, col) in enumerate(steps):
    x = 0.7 + i * 3.08
    sp = rect(s, x, 1.7, 2.85, 2.15, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    fill_text(sp, [(t, 14, True, SKY if col != GOLD else WHITE), (big, 30, True, WHITE), (d, 11.5, False, WHITE)])
    if i < 3:
        rect(s, x + 2.87, 2.55, 0.2, 0.45, GREY, shape=MSO_SHAPE.RIGHT_ARROW)
card(s, 0.7, 4.15, 5.9, 2.55, "Where WBI's value comes from",
     "• MM Studio — WBI's growth engine: $2.9M → $36.3M by 2030, the plan's biggest swing factor.\n"
     "• WBI core exports: $4.69M in 2026, +20%/yr, NTS-certified FY2025 base.\n"
     "• Technology upside (AI.pep, DPW licence) accrues inside GenApep — you participate via the block.",
     line=ACCENT, bsz=13.5)
card(s, 6.85, 4.15, 5.8, 2.55, "The protections that matter to you",
     "• A2: consideration only to shareholder bodies, pro rata — the draft's 40% to one individual "
     "(zero to other WBI shareholders) becomes impossible.\n"
     "• B4/B5: any Stage-2 deal carries every shareholder through as one block at one ratio.\n"
     "• A1: if Eyesel's real revenue departs >15% from $8.0M, equalisation is mandatory.",
     line=RED, bsz=13.5)

# ============================================================ 13 PRICE TEST
s = slide(); header(s, "The price test on First Vital", "What the draft's percentages imply — and how to check them")
rows = [
    ["GenApep combined value (planning)", "$30–40M", "JV plan valuation sheet (before AI.pep/DPW)"],
    ["GenApep share under the draft", "75%", "40% + 35% to the two individuals"],
    ["Implied combined group value", "$40–53M", "GenApep value ÷ 0.75"],
    ["Implied value of First Vital", "$10–13M", "the remaining 25%"],
]
table(s, ["Step", "Value", "Basis"], rows, 0.7, 1.6, 11.95, [4.6, 1.9, 5.45], fs=13, hdr_fs=13, hl_rows=(3,))
card(s, 0.7, 3.85, 11.95, 1.45, "The question to answer before any share moves",
     "Would First Vital be worth $10–13M if it were NOT public? If yes — on the GLP-1 business alone "
     "— the terms are fair and the listing status comes free. If not, most of the price pays for a "
     "filing record that has not yet been verified.", line=GOLD, tsz=16, bsz=14.5)
stat(s, 0.7, 5.5, 3.85, 1.2, "Gate 0", "CIK + full EDGAR index, in writing", fill=RED, bigsz=26)
stat(s, 4.75, 5.5, 3.85, 1.2, "This week", "request Jul–Aug GLP-1 actuals", fill=ACCENT, bigsz=24)
stat(s, 8.8, 5.5, 3.85, 1.2, "Then", "independent valuations + fairness opinion", fill=NAVY, bigsz=26)

# ============================================================ 14 PROTECTIONS
s = slide(); header(s, "Protections already drafted", "Sign Amendment No. 1 now; Addendum A when Route B is live")
card(s, 0.7, 1.6, 5.9, 5.1, "Amendment No. 1 (both scenarios)",
     "A1 · Equalisation becomes MANDATORY beyond 15% variance — independent valuer decides if "
     "boards cannot.\n"
     "A2 · Consideration to shareholder bodies pro rata, never to principals.\n"
     "A3/A4 · Third-Director interests disclosed; deadlocks resolved (pool fixed at 7.5% fallback).\n"
     "A5 · Technology chain of title = condition precedent to completion.\n"
     "A6/A7 · Technology encumbrances disclosed; provenance pack before sweat-equity vesting.\n"
     "A8 · No one may present GenApep as an existing rights-holding entity before it is one.\n"
     "A9 · Open items scheduled with owners and dates.", line=ACCENT, bsz=13.5)
card(s, 6.85, 1.6, 5.8, 5.1, "Addendum A (Route-B machinery)",
     "B1 · GenApep Korea is the ONLY permitted acquisition vehicle.\n"
     "B2 · Sweat pool carries into any new structure on equivalent terms.\n"
     "B3 · Route selection is a Reserved Matter; Route A stays the default.\n"
     "B4 · Ratio preservation: consideration strictly pro rata — binding.\n"
     "B5 · Block entry: one exchange ratio for the whole JV, never individual deals.\n"
     "B6/B11 · No unilateral negotiation; exclusivity extends to listing vehicles.\n"
     "B7 · Related-party deals need valuations, fairness opinion, disinterested approval.\n"
     "B8–B10, B12 · Verification gates, audit scoped for either route, consequences schedule.",
     line=GOLD, bsz=13.5)

# ============================================================ 15 TIMELINE
s = slide(); header(s, "Timeline & critical path", "The audit is the clock; everything else runs in parallel")
rows = [
    ["Now", "Sign Amendment No. 1 · request FVH EDGAR index + GLP-1 actuals", "Boards / FA"],
    ["≤ 30 days", "Eyesel updated financial report + principal confirmation", "Eyesel"],
    ["Q3–Q4 2026", "Incorporate GenApep Holdings (Cayman) + GenApep Korea", "FA / counsel"],
    ["Q4 2026", "Appoint PCAOB auditor — CRITICAL PATH · definitive SHA + pool KPIs", "Joint / FA"],
    ["Q1–Q2 2027", "Share exchange: GenApep Korea acquires WBI + Eyesel", "All parties"],
    ["When gated", "Route decision under Addendum A (verified FVH, valuations, fairness opinion)", "Boards"],
    ["≤ end-2027", "Confidential SEC draft registration (Route A) / combination filings (Route B)", "FA / US counsel"],
    ["2028", "US listing — market-dependent", "All"],
]
table(s, ["When", "Milestone", "Owner"], rows, 0.7, 1.6, 11.95, [1.7, 8.0, 2.25], fs=12.5, hdr_fs=13, hl_rows=(3,))
band = rect(s, 0.7, 5.75, 11.95, 0.95, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
fill_text(band, [("Two audited fiscal years to PCAOB standard are required before ANY filing — start the "
                  "auditor appointment now and the route choice steals no time.", 15, True, NAVY)])

# ============================================================ 16 DECISIONS
s = slide(); header(s, "Decisions requested today", "Approve these seven actions")
items = [
    ("1", "Sign Amendment No. 1 and start Stage 1 — both boards.", ACCENT),
    ("2", "Treat Addendum A as the gate for all First Vital discussion.", TEAL),
    ("3", "Instruct the FA to obtain FVH's CIK / EDGAR index and Jul–Aug GLP-1 actuals — this week.", GOLD),
    ("4", "Eyesel: updated financials + principal confirmation within 30 days.", DARK2),
    ("5", "Appoint the PCAOB auditor in Q4-2026, scoped for either route.", NAVY),
    ("6", "Do not sign the draft IPO MOU as written — revise to value-based block entry.", RED),
    ("7", "Progress the DPW licence to asset-level term sheets after chain-of-title is documented.", TEAL),
]
for i, (n, d, col) in enumerate(items):
    y = 1.6 + i * 0.76
    badge = rect(s, 0.8, y, 0.6, 0.6, col, shape=MSO_SHAPE.OVAL, shadow=True)
    fill_text(badge, [(n, 18, True, WHITE)])
    textbox(s, 1.65, y + 0.02, 11.0, 0.7, [(d, 15.5, i in (0, 5), NAVY)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.8, 7.0, 11.9, 0.4,
        [("PRIVATE & CONFIDENTIAL — WBI & Eyesel directors only · planning figures · not advice",
          11, False, GREY)])

prs.save(DIR + "GenApep_Executive_Presentation.pptx")
print("Saved GenApep_Executive_Presentation.pptx — %d slides" % len(prs.slides._sldIdLst))
