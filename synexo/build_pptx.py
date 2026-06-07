#!/usr/bin/env python3
"""Build the investor deck -> Synapep_Investor_Deck.pptx.

Graphic-driven, large-type deck. Numbers single-sourced from build_multimarket_model (RESULTS).
Body text >=18pt; rich vector graphics (flow, flywheel, hub-spoke, chevrons, timeline, charts).
"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import build_multimarket_model as mm

DIR = "/home/user/gamification/synexo/"
NAVY = RGBColor(0x1B, 0x33, 0x55); ACCENT = RGBColor(0x2E, 0x86, 0xC1)
TEAL = RGBColor(0x1A, 0xA1, 0x9C); GOLD = RGBColor(0xE0, 0x9E, 0x2B)
GREY = RGBColor(0x5A, 0x5A, 0x5A); LIGHT = RGBColor(0xEC, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); RED = RGBColor(0xC0, 0x39, 0x2B)
SKY = RGBColor(0xCF, 0xDD, 0xEA); DARK2 = RGBColor(0x24, 0x44, 0x6E)
Y = mm.YEARS; R = mm.RESULTS; ORDER = mm.ORDER

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
IN = Inches


CURRENT = [0]


def slide():
    CURRENT[0] += 1
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0, shadow=False):
    sp = s.shapes.add_shape(shape, IN(l), IN(t), IN(w), IN(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = shadow
    return sp


def fill_text(sp, lines, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    """lines: list of (text, size, bold, color)."""
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"


def textbox(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
        p.space_after = Pt(4)
    return tb


def header(s, title, n, kicker=None):
    rect(s, 0, 0, 13.333, 1.2, NAVY)
    rect(s, 0.55, 0.34, 0.14, 0.52, ACCENT)
    textbox(s, 0.85, 0.18, 11.6, 0.9, [(title, 27, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        textbox(s, 0.87, 0.86, 11.6, 0.3, [(kicker, 12, False, SKY)])
    textbox(s, 12.4, 7.04, 0.8, 0.35, [(str(CURRENT[0]), 12, False, GREY)], align=PP_ALIGN.RIGHT)


def bullets(s, items, l=0.9, t=1.6, w=11.6, h=5.2, size=19):
    tb = s.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("●  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size - lvl * 3); r.font.name = "Calibri"
        r.font.color.rgb = NAVY if lvl == 0 else GREY; r.font.bold = False
        p.space_after = Pt(12 if lvl == 0 else 6)
    return tb


def card(s, l, t, w, h, title, body, fill=WHITE, tcol=NAVY, bcol=GREY, tsz=18, bsz=15,
         line=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = rect(s, l, t, w, h, fill, line=line, line_w=1.5, shape=shape, shadow=True)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(10); tf.margin_right = Pt(10); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title; r.font.size = Pt(tsz); r.font.bold = True
    r.font.color.rgb = tcol; r.font.name = "Calibri"
    if body:
        p2 = tf.add_paragraph(); p2.space_before = Pt(4); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(bsz); r2.font.color.rgb = bcol
        r2.font.name = "Calibri"
    return sp


def stat(s, l, t, w, h, big, label, fill=NAVY, bigcol=WHITE, labcol=SKY, bigsz=36, labsz=14):
    sp = rect(s, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    fill_text(sp, [(big, bigsz, True, bigcol), (label, labsz, False, labcol)])
    return sp


def arrow(s, l, t, w, h, color=ACCENT, rotation=0, shape=MSO_SHAPE.RIGHT_ARROW):
    sp = rect(s, l, t, w, h, color, shape=shape)
    sp.rotation = rotation
    return sp


def chart_fonts(chart, size=14):
    try:
        chart.value_axis.tick_labels.font.size = Pt(size)
        chart.category_axis.tick_labels.font.size = Pt(size)
    except Exception:
        pass
    if chart.has_legend:
        chart.legend.font.size = Pt(size)


# =================================================================== 1 TITLE
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 13.333, 0.28, ACCENT)
rect(s, 0, 7.22, 13.333, 0.28, ACCENT)
# decorative circles
for (cx, cy, d, col) in [(11.4, 1.5, 2.6, DARK2), (12.2, 5.6, 1.8, DARK2), (10.6, 6.4, 1.0, ACCENT)]:
    o = rect(s, cx, cy, d, d, col, shape=MSO_SHAPE.OVAL)
rect(s, 0.95, 2.55, 0.2, 1.9, ACCENT)
textbox(s, 1.35, 2.45, 9.5, 1.3, [("SYNAPEP", 56, True, WHITE)])
textbox(s, 1.4, 3.85, 9.5, 0.9, [("AI-Personalised Aesthetic Medical Devices", 26, False, SKY)])
textbox(s, 1.4, 4.75, 9.5, 0.6, [("Korea-anchored, export-led  ·  Investor Deck", 17, False,
                                  RGBColor(0x9F, 0xB6, 0xCB))])
textbox(s, 1.4, 6.5, 10, 0.5, [("Illustrative — not investment advice. Figures are v0.4 scenario "
                                "estimates.", 12, False, RGBColor(0x86, 0x9C, 0xB4))])

# =================================================================== 2 OPPORTUNITY
s = slide(); header(s, "A large, fast-growing market", 2, "Medical aesthetics — global")
cd = CategoryChartData(); cd.categories = ["2024", "2033E"]
cd.add_series("Market size ($B)", (17.5, 56.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, IN(0.7), IN(1.7), IN(6.6), IN(5.2), cd)
ch = gf.chart; ch.has_legend = False
ch.value_axis.has_major_gridlines = False
pl = ch.plots[0]; pl.has_data_labels = True; pl.gap_width = 80
pl.data_labels.number_format = '"$"0"B"'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(18); pl.data_labels.font.bold = True; pl.data_labels.font.color.rgb = NAVY
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = ACCENT
chart_fonts(ch, 16)
stat(s, 8.0, 1.8, 4.6, 1.35, "~$56B", "global medical aesthetics by 2033", fill=NAVY)
stat(s, 8.0, 3.35, 2.2, 1.35, "10–13%", "CAGR", fill=TEAL, bigsz=30)
stat(s, 10.4, 3.35, 2.2, 1.35, ">40%", "injectables share", fill=GOLD, labcol=WHITE, bigsz=30)
card(s, 8.0, 4.9, 4.6, 1.75, "Why now",
     "Commodity injectables (HA, botox) keep growing but commoditise. Buyers want differentiated, "
     "regenerative, AI-personalised outcomes — a new lane, not a replacement.", bsz=15)

# =================================================================== 3 SOLUTION FLOW
s = slide(); header(s, "The solution — device + kits + AI", 3, "A razor-and-blades medical-device business")
y0 = 2.3; bw = 3.5; bh = 2.2
card(s, 0.7, y0, bw, bh, "1 · Applicator device",
     "Placed/sold to clinics. Also the data-capture endpoint. ~$3,000 ASP.", line=ACCENT, tsz=18, bsz=15)
arrow(s, 4.32, y0 + 0.8, 0.75, 0.6, ACCENT)
card(s, 5.15, y0, bw, bh, "2 · Personalised kits",
     "Recurring CodeLife-configured peptide/exosome treatment kits. ~$30–35 ex-factory, ~65% GM.",
     line=TEAL, tsz=18, bsz=15)
arrow(s, 8.77, y0 + 0.8, 0.75, 0.6, TEAL)
card(s, 9.6, y0, bw, bh, "3 · AI + data",
     "Every treatment feeds CodeLife — the model improves, the peptides improve, lock-in deepens.",
     line=GOLD, tsz=18, bsz=15)
rect(s, 0.7, 5.2, 11.93, 1.25, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 1.0, 5.35, 11.4, 1.0,
        [("Personalised by physician PURPOSE within an approved envelope — fast 'modified-device' "
          "registration, not bespoke per-patient devices.", 18, True, NAVY)], anchor=MSO_ANCHOR.MIDDLE)

# =================================================================== TWO CO-EQUAL MOATS
s = slide(); header(s, "Two co-equal moats", 0, "Either alone is copyable; together they compound")
card(s, 0.7, 1.65, 5.85, 4.45, "Moat A · AI engine + data flywheel",
     "CodeLife.AI designs peptides that outperform generalised products. An AI-QA suite validates "
     "every batch and feeds a compounding outcomes dataset — the durable asset. Physician "
     "decision-support within a cleared envelope (SaMD-aware).", line=ACCENT, tsz=20, bsz=17)
card(s, 6.8, 1.65, 5.85, 4.45, "Moat B · Physician co-development engine",
     "Elite Korean KOLs co-develop purpose-built SKUs. A fast/cheap regulatory 'modification "
     "engine' registers each one. A compliance-first benefit-sharing royalty makes the inventing "
     "doctor the champion. A relationship + process asset — hard to copy.", line=GOLD, tsz=20, bsz=17)
band = rect(s, 0.7, 6.25, 11.95, 0.75, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
fill_text(band, [("Devices + kits = the revenue wedge and the data-capture layer for Moat A", 16,
                  True, NAVY)])

# =================================================================== THREE ENGINES
s = slide(); header(s, "Synapep Labs — three engines", 0, "Moat B, in detail")
card(s, 0.7, 1.6, 3.85, 3.95, "1 · KOL co-development lab",
     "Korea's world-class physicians co-design products from the doctor's clinical perspective. "
     "Doctor-personalised, not patient-personalised. CodeLife accelerates each design.",
     line=ACCENT, tsz=18, bsz=15)
card(s, 4.75, 1.6, 3.85, 3.95, "2 · Modification engine",
     "Each variant is a minor change — Korea's improved-device / negative-list regime + same-active "
     "functional-cosmetic / quasi-drug routes. Capital-efficient SKU factory.", line=TEAL, tsz=18,
     bsz=15)
card(s, 8.8, 1.6, 3.85, 3.95, "3 · KOL royalty flywheel",
     "The inventing doctor champions the SKU to peers and shares economics as an FMV royalty for "
     "genuine IP co-invention — DECOUPLED from their own usage. K-Sunshine / anti-rebate compliant.",
     line=GOLD, tsz=18, bsz=15)
specs = [("<3%", "formulation Δ", TEAL), ("~3 mo", "approval cycle", ACCENT),
         ("<$30K", "per SKU", NAVY), ("~20", "SKUs / 6 mo", GOLD), ("<$600K", "total reg dev", DARK2)]
for i, (big, lab, col) in enumerate(specs):
    stat(s, 0.7 + i * 2.46, 5.8, 2.3, 1.05, big, lab, fill=col, bigsz=28, labsz=13)

# =================================================================== MOAT A · DATA FLYWHEEL
s = slide(); header(s, "Moat A — AI efficacy + data flywheel", 0,
                    "Performance, not merely personalisation")
cx, cy = 6.55, 4.35
o = rect(s, cx - 1.15, cy - 1.15, 2.3, 2.3, ACCENT, shape=MSO_SHAPE.OVAL, shadow=True)
fill_text(o, [("DATA", 20, True, WHITE), ("FLYWHEEL", 20, True, WHITE)])
nodes = [("Treatments\n& outcomes", 5.3, 1.7), ("Bigger\nconsented dataset", 9.6, 3.6),
         ("Better\nAI model", 6.4, 6.0), ("Better\npeptides", 2.4, 3.6)]
ncol = [TEAL, NAVY, GOLD, DARK2]
boxes = []
for (txt, l, t), col in zip(nodes, ncol):
    sp = rect(s, l, t, 2.5, 1.05, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    fill_text(sp, [(line, 15, True, WHITE) for line in [txt.replace("\n", " ")]])
    boxes.append((l + 1.25, t + 0.52))
# clockwise curved arrows around the hub
for (al, at, rot) in [(7.7, 2.5, 35), (8.0, 5.0, 125), (4.6, 5.2, 215), (4.3, 2.6, 305)]:
    a = rect(s, al, at, 1.0, 0.55, ACCENT, shape=MSO_SHAPE.RIGHT_ARROW); a.rotation = rot
textbox(s, 0.7, 6.75, 12, 0.5,
        [("Day-one asset: the consented dataset + data-rights + analytics. 'Own large model' is the "
          "direction, not the seed-stage claim.", 14, False, GREY)])

# =================================================================== 5 SAME PRODUCT / HUB-SPOKE
s = slide(); header(s, "Same product, more markets", 5, "One Korean-developed product, exported as-is")
hub_x, hub_y = 3.4, 4.4
markets = ["Hong Kong", "Vietnam", "Singapore", "Thailand", "Malaysia", "Indonesia",
           "Philippines", "Saudi", "UAE", "Taiwan", "China"]
# spokes (connectors first, behind)
import math as _m
positions = []
for i, mk in enumerate(markets):
    ang = -90 + i * (360.0 / len(markets))
    rad = _m.radians(ang)
    px = hub_x + 2.9 * _m.cos(rad); py = hub_y + 2.55 * _m.sin(rad)
    positions.append((px, py, mk))
for px, py, mk in positions:
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(hub_x), IN(hub_y), IN(px), IN(py))
    cn.line.color.rgb = RGBColor(0xC2, 0xD2, 0xE0); cn.line.width = Pt(1.5)
for px, py, mk in positions:
    chip = rect(s, px - 0.72, py - 0.28, 1.44, 0.56, LIGHT, line=ACCENT, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    fill_text(chip, [(mk, 13, True, NAVY)])
hub = rect(s, hub_x - 1.15, hub_y - 0.85, 2.3, 1.7, NAVY, shape=MSO_SHAPE.OVAL, shadow=True)
fill_text(hub, [("KOREA", 19, True, WHITE), ("MFDS", 19, True, WHITE),
                ("master dossier", 12, False, SKY)])
# right column points
card(s, 7.4, 1.7, 5.3, 1.25, "No per-market R&D",
     "Personalisation is by physician purpose within the SAME approved envelope.", bsz=15)
card(s, 7.4, 3.05, 5.3, 1.25, "Identical unit economics",
     "Only the clinic base and timing change country-to-country.", bsz=15, line=TEAL)
card(s, 7.4, 4.4, 5.3, 1.25, "Scale + data benefits",
     "One production / QC / stability spec; one comparable global dataset.", bsz=15, line=GOLD)
card(s, 7.4, 5.75, 5.3, 1.2, "Reliance pathway",
     "The Korean technical file is the master dossier for every export filing.", bsz=15, line=DARK2)

# =================================================================== 6 REGULATORY CHEVRON
s = slide(); header(s, "Regulatory strategy", 6, "Korea → Hong Kong → export, with the right product on the fast path")
steps = [("KOREA MFDS", "'Modified-device' route for the microneedling/topical kit + applicator (~months)", ACCENT),
         ("HONG KONG MDACS", "Light voluntary listing; MFDS is a recognised reference regulator", TEAL),
         ("EXPORT", "Reliance/recognition across ASEAN + GCC on the same Korean file", GOLD)]
for i, (t1, t2, col) in enumerate(steps):
    l = 0.7 + i * 4.15
    cv = rect(s, l, 2.0, 4.0, 1.7, col, shape=MSO_SHAPE.CHEVRON, shadow=True)
    fill_text(cv, [(t1, 19, True, WHITE), (" ", 6, False, WHITE), (t2, 13, False, WHITE)])
card(s, 0.7, 4.3, 12.0, 1.0, "Lead with the device family",
     "Microneedling/topical kit + applicator registers fast as a device. The IM/injectable line is "
     "likely a drug/biologic — run it on a separate, slower track so it never gates launch.",
     bsz=15, line=NAVY)
card(s, 0.7, 5.5, 12.0, 1.0, "CodeLife = configuration within a cleared envelope",
     "Physician in the loop; manages Software-as-a-Medical-Device exposure. Claim discipline on "
     "exosomes (non-human salmon/synthetic origin; function-based claims).", bsz=15, line=NAVY)

# =================================================================== 7 EXPORT TIERS
s = slide(); header(s, "Export springboard — KFDA reliance", 7, "Two tiers off one Korean registration")
def chips(s, items, l, t, w, per_row=2, cw=2.35, ch=0.62, gap=0.12, col=LIGHT):
    for i, it in enumerate(items):
        r = i // per_row; c = i % per_row
        x = l + c * (cw + gap); ycoord = t + r * (ch + gap)
        sp = rect(s, x, ycoord, cw, ch, col, line=ACCENT, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        fill_text(sp, [(it, 14, True, NAVY)])
p1 = rect(s, 0.7, 1.65, 6.0, 5.1, LIGHT, line=ACCENT, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 0.95, 1.8, 5.5, 0.8, [("TIER 1 · ASEAN + Hong Kong", 19, True, ACCENT),
                                 ("Fastest — reliance-led off the Korean file", 13, False, GREY)])
chips(s, ["Hong Kong", "Vietnam", "Singapore", "Thailand", "Malaysia", "Indonesia", "Philippines"],
      0.95, 3.05, 5.5, per_row=2)
textbox(s, 0.95, 6.1, 5.5, 0.6,
        [("HK MDACS reference · VN accepts MFDS · ASEAN CSDT cascade", 12, False, GREY)])
p2 = rect(s, 7.0, 1.65, 5.6, 5.1, RGBColor(0xEE, 0xE9, 0xDD), line=GOLD, line_w=1.5,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 7.25, 1.8, 5.2, 0.8, [("TIER 2 · GCC + Greater China", 19, True, GOLD),
                                 ("Large, attractive, more work per market", 13, False, GREY)])
chips(s, ["Saudi (SFDA)", "UAE (MOHAP)", "Taiwan", "China"], 7.25, 3.05, 5.2, per_row=2,
      col=RGBColor(0xF6, 0xF1, 0xE5))
textbox(s, 7.25, 6.1, 5.2, 0.6,
        [("SFDA regional reference · TFDA uses Korean data · China via Hainan", 12, False, GREY)])

# =================================================================== 8 FINANCIALS CHART
s = slide(); header(s, "Multi-market revenue — 3 scenarios", 8, "Same product; scenarios flex ramp, price, utilisation")
cd = CategoryChartData(); cd.categories = [str(y) for y in Y]
for scn in ORDER:
    cd.add_series(scn, tuple(R[scn][0][y]["tot_rev"] / 1e6 for y in Y))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, IN(0.7), IN(1.65), IN(8.0), IN(5.4), cd)
ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.value_axis.has_title = True; ch.value_axis.axis_title.text_frame.text = "Revenue ($M)"
ch.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.value_axis.has_major_gridlines = True
for i, col in enumerate([GREY, ACCENT, NAVY]):
    ch.series[i].format.fill.solid(); ch.series[i].format.fill.fore_color.rgb = col
chart_fonts(ch, 15)
for i, scn in enumerate(ORDER):
    r = R[scn][0][2030]; col = [GREY, ACCENT, NAVY][i]
    stat(s, 9.05, 1.75 + i * 1.72, 3.7, 1.5,
         "$%.0fM" % (r["tot_rev"] / 1e6),
         "%s · 2030 · %s clinics · $%.0fM EBITDA" % (scn, "{:,}".format(r["cum"]), r["ebitda"] / 1e6),
         fill=col, bigsz=34, labcol=WHITE if col != ACCENT else WHITE)

# =================================================================== 9 SCENARIO TABLE
s = slide(); header(s, "Scenario assumptions", 9, "The product never changes — only these levers")
rows = [["Markets", "11 (China excl.)", "12", "12"],
        ["Clinic ramp vs base", "0.62×", "1.00×", "1.30×"],
        ["Kit EXW", "$30.00", "$32.50", "$35.00"],
        ["Kit gross margin", "%.0f%%" % (mm.kit_gm(30) * 100), "%.0f%%" % (mm.kit_gm(32.5) * 100),
         "%.0f%%" % (mm.kit_gm(35) * 100)],
        ["Device GM", "35%", "40%", "42%"],
        ["Kits / clinic / yr", "1,200", "1,400", "1,500"]]
headers = ["Lever", "Conservative", "Base", "Upside"]
gt = s.shapes.add_table(len(rows) + 1, 4, IN(0.9), IN(1.8), IN(11.5), IN(4.6)).table
widths = [3.5, 2.7, 2.7, 2.6]
for i, w in enumerate(widths):
    gt.columns[i].width = IN(w)
for j, htext in enumerate(headers):
    c = gt.cell(0, j); c.text = htext; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE
for i, row in enumerate(rows, start=1):
    for j, v in enumerate(row):
        c = gt.cell(i, j); c.text = str(v)
        c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(16); r.font.color.rgb = NAVY; r.font.bold = (j == 0)

# =================================================================== 10 VALUATION
s = slide(); header(s, "Valuation — two lenses", 10, "Priced on the floor today; the AI upside is the Series A")
card(s, 0.7, 1.7, 5.7, 2.2, "Lens A — floor (priced today)",
     "Device / consumables business. $5–7M pre-money; $6M midpoint for the $1M kickoff. "
     "Recurring medical consumables + device + regulatory moat.", line=ACCENT, tsz=19, bsz=16)
card(s, 0.7, 4.05, 5.7, 2.2, "Lens B — upside (the prize)",
     "AI / data-science company. Comps shift medtech → AI/data multiples; reframes the 2028 "
     "Series A as a premium, institutional AI round.", line=GOLD, tsz=19, bsz=16)
cd = CategoryChartData(); cd.categories = ["WBI", "Eyesel", "Investor (SAFE)"]
cd.add_series("Post-raise", (60.0, 25.7, 14.3))
gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, IN(6.9), IN(1.7), IN(5.8), IN(4.9), cd)
ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.font.size = Pt(15)
pts = ch.plots[0].series[0].points
for pt, col in zip(pts, [NAVY, ACCENT, GOLD]):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
dl = ch.plots[0].data_labels; dl.number_format = '0.0"%"'; dl.number_format_is_linked = False
dl.font.size = Pt(15); dl.font.bold = True; dl.font.color.rgb = WHITE
ch.plots[0].has_data_labels = True
textbox(s, 6.9, 6.55, 5.8, 0.4, [("Cap table at $6.0M pre / $1.0M kickoff", 13, False, GREY)],
        align=PP_ALIGN.CENTER)

# =================================================================== 11 THE ASK
s = slide(); header(s, "The ask — milestone-gated SAFE", 11, "US$1.0M post-money SAFE · $6.0M cap · 20% discount · MFN")
t1 = rect(s, 0.7, 1.8, 5.7, 2.0, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
fill_text(t1, [("$400k", 40, True, WHITE), ("Tranche 1 — at close", 16, True, SKY),
               ("Conditions precedent (below)", 13, False, SKY)])
arrow(s, 6.55, 2.55, 0.7, 0.55, ACCENT)
t2 = rect(s, 7.35, 1.8, 5.3, 2.0, ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
fill_text(t2, [("$600k", 40, True, WHITE), ("Tranche 2 — on milestones", 16, True, WHITE),
               ("Released when proven (below)", 13, False, WHITE)])
card(s, 0.7, 4.0, 5.7, 2.7, "Conditions precedent",
     "•  WBI IP assigned/licensed into the JV\n•  Verify the $25M sample (entity/product/channel)\n"
     "•  Clinic data-rights clause\n•  Clean cap table", line=NAVY, tsz=18, bsz=16)
card(s, 7.35, 4.0, 5.3, 2.7, "Milestones (release Tranche 2)",
     "•  Cold-chain spec + COGS BOM\n•  5–10 anchor clinics with reorders\n"
     "•  One prospective clinical readout\n•  Regulatory classification memo", line=ACCENT, tsz=18, bsz=16)

# =================================================================== 12 ROADMAP TIMELINE
s = slide(); header(s, "Roadmap", 12, "Korea anchor → reliance-led export waves")
rect(s, 0.9, 3.95, 11.5, 0.16, RGBColor(0xC2, 0xD2, 0xE0))
phases = [("2027", "Anchor", "Korea registration; anchor clinics; CodeLife v1", ACCENT, True),
          ("2027–28", "Wave 1", "Hong Kong · Vietnam · Singapore; reorder evidence", TEAL, False),
          ("2028–29", "Wave 2", "ASEAN cascade + GCC (Saudi, UAE); Series A", GOLD, True),
          ("2029–30", "Wave 3", "Taiwan · China (Hainan→NMPA)", DARK2, False)]
n = len(phases)
for i, (yr, ph, desc, col, up) in enumerate(phases):
    cx = 1.9 + i * 3.25
    dot = rect(s, cx - 0.28, 3.74, 0.56, 0.56, col, shape=MSO_SHAPE.OVAL, shadow=True)
    fill_text(dot, [(str(i + 1), 18, True, WHITE)])
    if up:
        cardtop = 1.75
    else:
        cardtop = 4.55
    cd2 = rect(s, cx - 1.45, cardtop, 2.9, 1.6, WHITE, line=col, line_w=2.0,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    tf = cd2.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    fill_text(cd2, [(yr + " · " + ph, 16, True, col), (desc, 14, False, GREY)])
    # connector dot->card
    cy2 = cardtop + (1.6 if up else 0)
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(cx), IN(4.03), IN(cx), IN(cy2))
    cn.line.color.rgb = col; cn.line.width = Pt(1.5)

# =================================================================== 13 WHAT WE'LL PROVE
s = slide(); header(s, "What we will prove (diligence-ready)", 13, "Turning the critique into milestones")
items = [("RED", "Cold chain & COGS", "Reformulate to 2–8 °C / ambient; publish a bottom-up BOM", RED),
         ("1", "Validation pack", "The $25M at entity/product/channel level + reorder cohorts", ACCENT),
         ("2", "Formulation envelope", "Clean AND clinically differentiating within the family", TEAL),
         ("3", "Exosome discipline", "Non-human (salmon/synthetic) origin; function-based claims", GOLD),
         ("4", "Data moat", "CodeLife v1 demo + data-rights in the clinic contract", DARK2),
         ("5", "Anchor evidence", "5–10 clinics live with documented reorders", NAVY)]
for i, (tag, title, body, col) in enumerate(items):
    r = i // 2; c = i % 2
    x = 0.7 + c * 6.15; ytop = 1.7 + r * 1.72
    badge = rect(s, x, ytop, 1.0, 1.0, col, shape=MSO_SHAPE.OVAL, shadow=True)
    fill_text(badge, [(tag, 16 if len(tag) > 1 else 22, True, WHITE)])
    tb = rect(s, x + 1.15, ytop, 4.75, 1.45, WHITE, line=col, line_w=1.5,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]; rr = p.add_run(); rr.text = title; rr.font.size = Pt(17); rr.font.bold = True
    rr.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(14)
    r2.font.color.rgb = GREY

# =================================================================== 14 CLOSING
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 13.333, 0.28, ACCENT); rect(s, 0, 7.22, 13.333, 0.28, ACCENT)
rect(s, 11.0, 1.2, 2.2, 2.2, DARK2, shape=MSO_SHAPE.OVAL)
rect(s, 11.8, 4.6, 1.4, 1.4, ACCENT, shape=MSO_SHAPE.OVAL)
rect(s, 0.95, 2.6, 0.2, 1.6, ACCENT)
textbox(s, 1.35, 2.5, 10.5, 1.2, [("Same product. More markets. A data moat.", 36, True, WHITE)])
textbox(s, 1.4, 3.95, 10.8, 1.6,
        [("Base case ~$%.0fM revenue / ~$%.0fM EBITDA by 2030 across 12 markets on one "
          "Korean-developed product." % (R["Base"][0][2030]["tot_rev"] / 1e6,
                                         R["Base"][0][2030]["ebitda"] / 1e6), 20, False, SKY)])
textbox(s, 1.4, 5.5, 10.8, 0.7, [("$1.0M milestone-gated kickoff · $6.0M cap · 20% discount", 18,
                                  True, WHITE)])

prs.save(DIR + "Synapep_Investor_Deck.pptx")
print("Saved Synapep_Investor_Deck.pptx — %d slides" % len(prs.slides._sldIdLst))
