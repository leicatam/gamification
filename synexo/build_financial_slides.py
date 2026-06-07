#!/usr/bin/env python3
"""Build a 3-slide financial summary -> Synapep_Financial_Slides.pptx.

Three slides off the bottom-up operating model: (1) Assumptions, (2) Baseline P&L,
(3) ROI / returns. Numbers single-sourced from build_financials (P, pnl, YEARS, ...).
Matches the investor-deck look (16:9, navy/accent palette, large type).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import build_financials as fin

DIR = "/home/user/gamification/synexo/"
NAVY = RGBColor(0x1B, 0x33, 0x55); ACCENT = RGBColor(0x2E, 0x86, 0xC1)
TEAL = RGBColor(0x1A, 0xA1, 0x9C); GOLD = RGBColor(0xE0, 0x9E, 0x2B)
GREY = RGBColor(0x5A, 0x5A, 0x5A); LIGHT = RGBColor(0xEC, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); RED = RGBColor(0xC0, 0x39, 0x2B)
SKY = RGBColor(0xCF, 0xDD, 0xEA); DARK2 = RGBColor(0x24, 0x44, 0x6E)

Y = fin.YEARS; P = fin.P
m = fin.m  # dollars -> $M

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
IN = Inches
CURRENT = [0]


def slide():
    CURRENT[0] += 1
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0, shadow=False):
    sp = s.shapes.add_shape(shape, IN(l), IN(t), IN(w), IN(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = shadow
    return sp


def fill_text(sp, lines, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"


def textbox(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
        p.space_after = Pt(4)
    return tb


def header(s, title, kicker=None):
    rect(s, 0, 0, 13.333, 1.2, NAVY)
    rect(s, 0.55, 0.34, 0.14, 0.52, ACCENT)
    textbox(s, 0.85, 0.18, 11.6, 0.9, [(title, 27, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        textbox(s, 0.87, 0.86, 11.6, 0.3, [(kicker, 14, False, SKY)])
    textbox(s, 12.4, 7.04, 0.8, 0.35, [("A%d" % CURRENT[0], 12, False, GREY)], align=PP_ALIGN.RIGHT)


def card(s, l, t, w, h, title, body, fill=WHITE, tcol=NAVY, bcol=GREY, tsz=18, bsz=16,
         line=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = rect(s, l, t, w, h, fill, line=line, line_w=1.5, shape=shape, shadow=True)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(10); tf.margin_right = Pt(10); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title; r.font.size = Pt(tsz); r.font.bold = True
    r.font.color.rgb = tcol; r.font.name = "Calibri"
    if body:
        p2 = tf.add_paragraph(); p2.space_before = Pt(4); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(bsz); r2.font.color.rgb = bcol
        r2.font.name = "Calibri"
    return sp


def stat(s, l, t, w, h, big, label, fill=NAVY, bigcol=WHITE, labcol=SKY, bigsz=36, labsz=14):
    sp = rect(s, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    fill_text(sp, [(big, bigsz, True, bigcol), (label, labsz, False, labcol)])
    return sp


def chart_fonts(chart, size=14):
    try:
        chart.value_axis.tick_labels.font.size = Pt(size)
        chart.category_axis.tick_labels.font.size = Pt(size)
    except Exception:
        pass
    if chart.has_legend:
        chart.legend.font.size = Pt(size)


# ---- derived figures (single-sourced) -------------------------------------
GP_VIAL = fin.GP_VIAL; PRICE = fin.PRICE; COST = fin.COST
GM = P[2030]["gm"]
cum_rev = sum(P[y]["rev"] for y in Y)
cum_eb = sum(P[y]["ebitda"] for y in Y)
eb2030 = P[2030]["ebitda"]
KICKOFF = 1.0e6           # $1.0M kickoff SAFE (deck terms)
CAP_POST = 7.0e6          # $6.0M pre + $1.0M = $7.0M post
OWN = KICKOFF / CAP_POST  # investor post-money ownership ≈ 14.3%

# cumulative-EBITDA payback of the kickoff
cum = 0.0; payback_year = None
for y in Y:
    cum += P[y]["ebitda"]
    if payback_year is None and cum >= KICKOFF:
        payback_year = y

# illustrative equity ROI: exit EV = multiple × 2030 EBITDA, investor takes OWN
MULTS = [6, 8, 10]


def investor_value(mult):
    return OWN * mult * eb2030


# =================================================================== 1 ASSUMPTIONS
s = slide(); header(s, "Financial model — assumptions",
                    "Bottom-up vial (razor-and-blades) consumable model · from Jul 2026, production Feb 2027")
# left: input cards
rows = [
    ("Ex-factory price / vial", "$%.0f–%.0f  (base $%.2f)" % (fin.PRICE_LO, fin.PRICE_HI, PRICE), ACCENT),
    ("Cost / vial", "$%.0f  =  bottoming $%.0f + production $%.0f" % (COST, fin.COST_EYESEL, fin.COST_PROD), TEAL),
    ("Gross profit / vial", "$%.2f   (%.1f%% margin)" % (GP_VIAL, GP_VIAL / PRICE * 100), GOLD),
    ("Vials / clinic / month", "%d   (MOQ %d per order)" % (fin.VIALS_PER_CLINIC_MO, fin.MOQ), DARK2),
    ("Production capacity", "%s vials/mo  (≈ %d clinics) — 4 lab + 2 shared" %
     (f"{fin.CAP_PER_MO:,}", fin.CAP_PER_MO // fin.VIALS_PER_CLINIC_MO), NAVY),
]
ytop = 1.55
for i, (lab, val, col) in enumerate(rows):
    yc = ytop + i * 0.92
    rect(s, 0.7, yc, 0.14, 0.74, col)
    rect(s, 0.84, yc, 6.7, 0.74, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, 1.05, yc + 0.04, 6.4, 0.68,
            [(lab, 15, True, NAVY), (val, 14, False, GREY)], anchor=MSO_ANCHOR.MIDDLE)

# right: clinic ramp assumption (the volume driver) + headline unit economics
stat(s, 7.85, 1.55, 2.35, 1.4, "$%.2f" % GP_VIAL, "gross profit / vial", fill=GOLD, labcol=WHITE, bigsz=34)
stat(s, 10.3, 1.55, 2.35, 1.4, "%.1f%%" % (GP_VIAL / PRICE * 100), "gross margin", fill=TEAL, bigsz=34)
card(s, 7.85, 3.15, 4.8, 1.35, "Volume driver — active clinics",
     "End-2027 110  →  2028 220  →  2029 350  →  2030 480.  100 vials/clinic/mo. "
     "2028–30 counts are assumptions to confirm.", line=ACCENT, tsz=16, bsz=14)
# mini clinic-count bars
cd = CategoryChartData(); cd.categories = [str(y) for y in Y]
cd.add_series("Avg active clinics", tuple(fin.avg_active(y) for y in Y))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, IN(7.85), IN(4.65), IN(4.8), IN(2.25), cd)
ch = gf.chart; ch.has_legend = False; ch.value_axis.has_major_gridlines = False
pl = ch.plots[0]; pl.has_data_labels = True; pl.gap_width = 60
pl.data_labels.number_format = '0'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(12); pl.data_labels.font.bold = True; pl.data_labels.font.color.rgb = NAVY
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = ACCENT
chart_fonts(ch, 12)
textbox(s, 7.85, 6.95, 4.8, 0.35, [("Avg active clinics / year — the model's revenue driver", 12, False, GREY)])

# =================================================================== 2 BASELINE
s = slide(); header(s, "Baseline — 5-year P&L",
                    "Base case (ex-factory $%.2f / vial) · USD · illustrative" % PRICE)
# revenue + EBITDA combo
cd = CategoryChartData(); cd.categories = [str(y) for y in Y]
cd.add_series("Revenue", tuple(m(P[y]["rev"]) for y in Y))
cd.add_series("EBITDA", tuple(m(P[y]["ebitda"]) for y in Y))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, IN(0.7), IN(1.55), IN(7.6), IN(4.05), cd)
ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.value_axis.has_title = True; ch.value_axis.axis_title.text_frame.text = "$M"
ch.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
pl = ch.plots[0]; pl.gap_width = 90; pl.has_data_labels = True
pl.data_labels.number_format = '"$"0.0'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(11); pl.data_labels.font.bold = True
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = ACCENT
ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb = NAVY
chart_fonts(ch, 13)
# headline stats column
stat(s, 8.55, 1.55, 4.1, 1.2, "$%.1fM" % m(P[2030]["rev"]), "2030 revenue", fill=ACCENT, bigsz=32)
stat(s, 8.55, 2.85, 4.1, 1.2, "$%.1fM" % m(eb2030), "2030 EBITDA  (%.0f%% margin)" %
     (eb2030 / P[2030]["rev"] * 100), fill=NAVY, bigsz=32)
stat(s, 8.55, 4.15, 2.0, 1.45, "%.1f%%" % (GM * 100), "gross margin", fill=TEAL, bigsz=28)
stat(s, 10.65, 4.15, 2.0, 1.45, "2027", "EBITDA-positive", fill=GOLD, labcol=WHITE, bigsz=28)
# P&L strip table along the bottom
labels = ["$M", "Revenue", "Gross profit", "EBITDA"]
gt = s.shapes.add_table(4, 6, IN(0.7), IN(5.85), IN(11.95), IN(1.45)).table
gt.columns[0].width = IN(2.75)
for j in range(5):
    gt.columns[j + 1].width = IN(1.84)
hdr = ["$M"] + [str(y) for y in Y]
data = [
    ["Revenue"] + ["%.2f" % m(P[y]["rev"]) for y in Y],
    ["Gross profit"] + ["%.2f" % m(P[y]["gp"]) for y in Y],
    ["EBITDA"] + ["%.2f" % m(P[y]["ebitda"]) for y in Y],
]
for j, h in enumerate(hdr):
    c = gt.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
        for r in p.runs: r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
for i, row in enumerate(data, start=1):
    for j, v in enumerate(row):
        c = gt.cell(i, j); c.text = str(v)
        c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(13); r.font.color.rgb = NAVY; r.font.bold = (j == 0 or row[0] == "EBITDA")

# =================================================================== 3 ROI
s = slide(); header(s, "Return on investment",
                    "$1.0M milestone-gated kickoff · self-funding from first production · illustrative")
# top: capital-efficiency stats
stat(s, 0.7, 1.55, 3.85, 1.5, "$%.1fM" % m(cum_eb), "5-yr cumulative EBITDA", fill=NAVY, bigsz=34)
stat(s, 4.7, 1.55, 3.85, 1.5, "%.0f×" % (cum_eb / KICKOFF), "5-yr EBITDA / $1.0M in", fill=ACCENT, bigsz=34)
stat(s, 8.7, 1.55, 3.95, 1.5, "%d" % payback_year, "kickoff repaid by cum. EBITDA", fill=TEAL, bigsz=34)
# left: why capital-efficient
card(s, 0.7, 3.25, 5.85, 3.55, "Why so capital-efficient",
     "•  EBITDA-positive from 2027 — the first production year.\n"
     "•  %.1f%% gross margin; the 4-lab + 2-shared team is fixed to 60K vials/mo, so margin "
     "expands as clinics grow.\n"
     "•  Cumulative EBITDA covers the $1.0M kickoff by %d — no large survival round.\n"
     "•  $1.0M funds prep + working capital to self-sustaining." % (GM * 100, payback_year),
     line=ACCENT, tsz=18, bsz=15)
# right: illustrative equity ROI at exit
card(s, 6.7, 3.25, 5.95, 1.05, "Illustrative equity return (exit on 2030 EBITDA)",
     "Investor owns ~%.1f%% ($1.0M at $6.0M pre / $7.0M post). Exit EV = multiple × $%.1fM 2030 EBITDA." %
     (OWN * 100, m(eb2030)), line=GOLD, tsz=16, bsz=14)
heads = ["Exit multiple", "Enterprise value", "Investor stake", "Return on $1.0M"]
gt = s.shapes.add_table(len(MULTS) + 1, 4, IN(6.7), IN(4.45), IN(5.95), IN(2.35)).table
for i, w in enumerate([1.45, 1.7, 1.5, 1.3]):
    gt.columns[i].width = IN(w)
for j, h in enumerate(heads):
    c = gt.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = WHITE
for i, mult in enumerate(MULTS, start=1):
    ev = mult * eb2030; iv = investor_value(mult)
    vals = ["%d×" % mult, "$%.0fM" % m(ev), "$%.1fM" % m(iv), "~%.0f×" % (iv / KICKOFF)]
    for j, v in enumerate(vals):
        c = gt.cell(i, j); c.text = v
        c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(14); r.font.color.rgb = NAVY; r.font.bold = (j == 3)
textbox(s, 0.7, 6.95, 12.0, 0.4,
        [("Illustrative only — not investment advice. Equity ROI excludes dilution from future rounds; "
          "exit multiples are scenario placeholders.", 11, False, GREY)])

prs.save(DIR + "Synapep_Financial_Slides.pptx")
print("Saved Synapep_Financial_Slides.pptx — %d slides" % len(prs.slides._sldIdLst))
print("cum EBITDA $%.2fM | 2030 EBITDA $%.2fM | own %.1f%% | payback %s" %
      (m(cum_eb), m(eb2030), OWN * 100, payback_year))
