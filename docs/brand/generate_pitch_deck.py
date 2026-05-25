"""Generate the OPC Tutor investor pitch deck (.pptx).

Run:
    pip install python-pptx
    python docs/brand/generate_pitch_deck.py

Outputs:
    docs/brand/pitch-deck.pptx

The palette, type and voice come from docs/brand/brand-manual.md. Edit the
SLIDES list to change the narrative; edit the constants below to re-skin.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------- brand tokens ----------

INK_BLUE = RGBColor(0x0E, 0x2A, 0x47)
JADE = RGBColor(0x2E, 0x7D, 0x6B)
AMBER = RGBColor(0xD8, 0x9B, 0x3A)
BONE = RGBColor(0xF6, 0xF1, 0xE7)
CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
MIST = RGBColor(0xE4, 0xE9, 0xEE)
BRICK = RGBColor(0xA8, 0x41, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEADLINE_FONT = "IBM Plex Serif"
BODY_FONT = "Inter"
MONO_FONT = "IBM Plex Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN * 2


# ---------- helpers ----------


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color: RGBColor):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(rect, color)
    return rect


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    size: int = 14,
    font: str = BODY_FONT,
    color: RGBColor = CHARCOAL,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items: list[str],
    *,
    size: int = 16,
    color: RGBColor = CHARCOAL,
    bullet_color: RGBColor = JADE,
    spacing: int = 8,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        bullet = p.add_run()
        bullet.text = "●  "
        bullet.font.name = BODY_FONT
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        bullet.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_footer(slide, slide_no: int, total: int, label: str = "OPC Tutor — Investor brief"):
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), INK_BLUE)
    add_text(
        slide,
        MARGIN,
        SLIDE_H - Inches(0.32),
        CONTENT_W - Inches(1),
        Inches(0.3),
        label,
        size=10,
        color=BONE,
    )
    add_text(
        slide,
        SLIDE_W - MARGIN - Inches(1),
        SLIDE_H - Inches(0.32),
        Inches(1),
        Inches(0.3),
        f"{slide_no:02d} / {total:02d}",
        size=10,
        color=BONE,
        align=PP_ALIGN.RIGHT,
        font=MONO_FONT,
    )


def add_brand_bar(slide):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), AMBER)


def add_section_header(slide, eyebrow: str, headline: str):
    add_brand_bar(slide)
    add_text(
        slide,
        MARGIN,
        Inches(0.45),
        CONTENT_W,
        Inches(0.35),
        eyebrow.upper(),
        size=11,
        color=JADE,
        bold=True,
        font=MONO_FONT,
    )
    add_text(
        slide,
        MARGIN,
        Inches(0.85),
        CONTENT_W,
        Inches(1.4),
        headline,
        size=32,
        color=INK_BLUE,
        bold=True,
        font=HEADLINE_FONT,
    )


# ---------- slide builders ----------


def slide_title(prs, n=None, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill(slide.background.fill, BONE) if False else None  # background handled by full rect
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, INK_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), AMBER)

    # Wordmark proxy
    add_text(
        slide,
        Inches(1.0),
        Inches(0.9),
        Inches(5),
        Inches(0.6),
        "OPC TUTOR",
        size=22,
        color=INK_BLUE,
        bold=True,
        font=HEADLINE_FONT,
    )
    # the dot above the P
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.62), Inches(0.78), Inches(0.16), Inches(0.16)
    )
    set_fill(dot, AMBER)

    add_text(
        slide,
        Inches(1.0),
        Inches(2.6),
        Inches(11),
        Inches(2.4),
        "From learner\nto operator.",
        size=66,
        color=INK_BLUE,
        bold=True,
        font=HEADLINE_FONT,
    )
    add_text(
        slide,
        Inches(1.0),
        Inches(5.0),
        Inches(11),
        Inches(0.6),
        "Adaptive AI tutor and agent toolkit for the Hong Kong one-person Ltd.",
        size=22,
        color=CHARCOAL,
    )
    add_text(
        slide,
        Inches(1.0),
        Inches(6.6),
        Inches(11),
        Inches(0.4),
        "Investor brief  ·  Phase 1 (Arc 1) shipping  ·  Hong Kong, Cap. 622",
        size=12,
        color=JADE,
        font=MONO_FONT,
    )


def slide_pitch(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "The 30-second pitch", "One human, one company, one agent stack.")

    add_text(
        slide,
        MARGIN,
        Inches(2.4),
        CONTENT_W,
        Inches(3.5),
        (
            "Hong Kong is the first super-aged Asian city with a clean one-person\n"
            "private-company framework and a two-tier profits tax that rewards small\n"
            "incorporation. Agentic AI has collapsed the cost of running a company\n"
            "without collapsing the cost of deciding what it does.\n\n"
            "OPC Tutor takes a Hong Kong professional from “what is an OPC” to\n"
            "operating one — adaptively, in six arcs — then ships the draft-only\n"
            "agent stack that runs the books, the calendar, and the paperwork."
        ),
        size=20,
        color=CHARCOAL,
    )
    add_footer(slide, n, total)


def slide_two_waves(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Why now", "Two waves are meeting in Hong Kong.")

    col_w = (CONTENT_W - Inches(0.4)) / 2

    def card(x, title, body, accent):
        add_rect(slide, x, Inches(2.4), col_w, Inches(4.2), WHITE)
        add_rect(slide, x, Inches(2.4), Inches(0.12), Inches(4.2), accent)
        add_text(
            slide,
            x + Inches(0.4),
            Inches(2.6),
            col_w - Inches(0.6),
            Inches(0.6),
            title,
            size=22,
            color=INK_BLUE,
            bold=True,
            font=HEADLINE_FONT,
        )
        add_text(
            slide,
            x + Inches(0.4),
            Inches(3.3),
            col_w - Inches(0.6),
            Inches(3.0),
            body,
            size=15,
            color=CHARCOAL,
        )

    card(
        MARGIN,
        "Aging society",
        (
            "Hong Kong crossed the super-aged line (>20% over 65) around 2024, on\n"
            "track for ~30% by the mid-2030s. MPF replaces only a fraction of\n"
            "pre-retirement income for most professionals.\n\n"
            "A large cohort holds 20–35 years of expensive expertise the city\n"
            "needs to keep productive."
        ),
        JADE,
    )
    card(
        MARGIN + col_w + Inches(0.4),
        "Agentic AI",
        (
            "A company is mostly a coordination structure: books, billing,\n"
            "drafting, calendar, filings. Agents now perform draft-quality\n"
            "versions of all of it.\n\n"
            "The smallest unit of viable incorporation has shrunk. One human,\n"
            "plus a stack, can do what used to need a small firm."
        ),
        AMBER,
    )
    add_footer(slide, n, total)


def slide_hk_numbers(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Hong Kong", "The numbers that frame the wedge.")

    stats = [
        ("20%+", "of HK residents are over 65 today",
            "Crossed the super-aged threshold around 2024."),
        ("8.25%", "profits tax on the first HK$2M",
            "Two-tier rate; 16.5% above. Structurally favours small Ltd."),
        ("1", "director and 1 shareholder permitted",
            "Cap. 622 explicitly allows a single-person private Ltd."),
        ("42 days", "to file the NAR1 each year",
            "The product tracks every CR / IRD / MPFA deadline."),
    ]

    col_w = (CONTENT_W - Inches(0.6)) / 2
    row_h = Inches(1.9)
    for i, (big, label, sub) in enumerate(stats):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (col_w + Inches(0.6))
        y = Inches(2.4) + row * (row_h + Inches(0.3))
        add_rect(slide, x, y, col_w, row_h, WHITE)
        add_text(slide, x + Inches(0.4), y + Inches(0.15), col_w - Inches(0.6), Inches(0.9),
                 big, size=44, color=INK_BLUE, bold=True, font=HEADLINE_FONT)
        add_text(slide, x + Inches(0.4), y + Inches(0.95), col_w - Inches(0.6), Inches(0.4),
                 label, size=13, color=JADE, bold=True, font=MONO_FONT)
        add_text(slide, x + Inches(0.4), y + Inches(1.3), col_w - Inches(0.6), Inches(0.6),
                 sub, size=12, color=CHARCOAL)

    add_footer(slide, n, total)


def slide_unit(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Thesis", "The unit of viable incorporation just shrank.")
    add_text(
        slide, MARGIN, Inches(2.4), CONTENT_W, Inches(0.6),
        "Before",
        size=14, color=JADE, bold=True, font=MONO_FONT,
    )
    add_text(
        slide, MARGIN, Inches(2.9), CONTENT_W, Inches(1.0),
        "A small firm = 3–8 people: a partner, an accountant, a secretary,\n"
        "a paralegal, a junior or two. Overhead and coordination set a floor.",
        size=18, color=CHARCOAL,
    )
    add_text(
        slide, MARGIN, Inches(4.3), CONTENT_W, Inches(0.6),
        "Now",
        size=14, color=AMBER, bold=True, font=MONO_FONT,
    )
    add_text(
        slide, MARGIN, Inches(4.8), CONTENT_W, Inches(1.8),
        "One human + a draft-only agent stack: books, billing, calendar,\n"
        "research, filings drafted. Human decides, reviews, approves.\n"
        "The floor moved. The opportunity is to teach the human to use it.",
        size=18, color=CHARCOAL,
    )
    add_footer(slide, n, total)


def slide_what_is_opc(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "The OPC", "A Hong Kong private Ltd, one human, one agent stack.")

    headers = ["Dimension", "Sole proprietor (BR only)", "HK private Ltd (Cap. 622)"]
    rows = [
        ("Legal personality", "None — you are the business", "Separate legal person"),
        ("Liability", "Personal, unlimited", "Limited to share capital + director duties"),
        ("Tax", "Personal salaries-tax bands", "8.25% on first HK$2M, 16.5% above"),
        ("Compliance", "BR + IRD individual return", "BR + NNC1/NAR1/SCR + BIR51 + audit"),
        ("Succession", "Ends with you", "Shares transferable; company continues"),
        ("Perception", "Freelancer", "Company — matters for some clients"),
    ]

    top = Inches(2.4)
    row_h = Inches(0.45)
    col_widths = [Inches(2.6), Inches(4.6), Inches(4.9)]
    x_offsets = [MARGIN]
    for w in col_widths[:-1]:
        x_offsets.append(x_offsets[-1] + w)

    # header row
    for i, h in enumerate(headers):
        add_rect(slide, x_offsets[i], top, col_widths[i], row_h, INK_BLUE)
        add_text(
            slide, x_offsets[i] + Inches(0.15), top + Inches(0.05),
            col_widths[i] - Inches(0.2), row_h,
            h, size=12, color=BONE, bold=True, font=MONO_FONT,
        )

    # body
    for r, row in enumerate(rows):
        bg = WHITE if r % 2 == 0 else MIST
        y = top + row_h * (r + 1)
        for i, cell in enumerate(row):
            add_rect(slide, x_offsets[i], y, col_widths[i], row_h, bg)
            color = INK_BLUE if i == 0 else CHARCOAL
            bold = i == 0
            add_text(
                slide, x_offsets[i] + Inches(0.15), y + Inches(0.05),
                col_widths[i] - Inches(0.2), row_h,
                cell, size=12, color=color, bold=bold,
            )

    add_footer(slide, n, total)


def slide_product(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Product", "Six arcs. One adaptive tutor. A draft-only toolkit.")

    arcs = [
        ("01", "Why an OPC, why now", "Shipped"),
        ("02", "Incorporating in HK", "Outline"),
        ("03", "Running it", "Outline"),
        ("04", "The agent stack", "Outline"),
        ("05", "Risk, liability, insurance", "Outline"),
        ("06", "Succession & exit", "Outline"),
    ]

    col_w = (CONTENT_W - Inches(0.4) * 2) / 3
    row_h = Inches(1.7)
    for i, (num, title, status) in enumerate(arcs):
        col = i % 3
        row = i // 3
        x = MARGIN + col * (col_w + Inches(0.4))
        y = Inches(2.4) + row * (row_h + Inches(0.3))
        add_rect(slide, x, y, col_w, row_h, WHITE)
        add_rect(slide, x, y, Inches(0.12), row_h, INK_BLUE)
        add_text(slide, x + Inches(0.35), y + Inches(0.2), col_w, Inches(0.5),
                 f"Arc {num}", size=12, color=JADE, bold=True, font=MONO_FONT)
        add_text(slide, x + Inches(0.35), y + Inches(0.55), col_w - Inches(0.4), Inches(0.7),
                 title, size=18, color=INK_BLUE, bold=True, font=HEADLINE_FONT)
        chip_color = AMBER if status == "Shipped" else MIST
        chip_text_color = BONE if status == "Shipped" else CHARCOAL
        chip = add_rect(slide, x + Inches(0.35), y + Inches(1.2), Inches(1.2), Inches(0.32), chip_color)
        add_text(slide, x + Inches(0.4), y + Inches(1.24), Inches(1.1), Inches(0.3),
                 status.upper(), size=10, color=chip_text_color, bold=True, font=MONO_FONT,
                 align=PP_ALIGN.CENTER)

    add_footer(slide, n, total)


def slide_personas(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Adapts to the learner", "Three personas, three pacing defaults.")

    personas = [
        (
            "Mid-career pivot",
            "35–50",
            "AI-native, time-poor, mobile-first.",
            ["Terse, peer-to-peer tone", "Brisk pacing", "Wants tax + agent stack"],
            JADE,
        ),
        (
            "Pre-retirement professional",
            "50–65",
            "20–35 years of expertise. Skeptical of hype.",
            ["Substantive paragraphs", "Measured pacing", "Wants depth on tax + audit"],
            INK_BLUE,
        ),
        (
            "Active retiree",
            "65+",
            "Accessibility-first. Often referred by family.",
            ["Warm, plain English", "Voice I/O on, extra-large type", "Human escalation always offered"],
            AMBER,
        ),
    ]

    col_w = (CONTENT_W - Inches(0.4) * 2) / 3
    for i, (name, age, blurb, items, accent) in enumerate(personas):
        x = MARGIN + i * (col_w + Inches(0.4))
        y = Inches(2.4)
        h = Inches(4.4)
        add_rect(slide, x, y, col_w, h, WHITE)
        add_rect(slide, x, y, col_w, Inches(0.5), accent)
        add_text(slide, x + Inches(0.3), y + Inches(0.08), col_w - Inches(0.4), Inches(0.4),
                 name.upper(), size=12, color=BONE, bold=True, font=MONO_FONT)
        add_text(slide, x + Inches(0.3), y + Inches(0.7), col_w - Inches(0.4), Inches(0.4),
                 f"Age {age}", size=12, color=JADE, font=MONO_FONT, bold=True)
        add_text(slide, x + Inches(0.3), y + Inches(1.1), col_w - Inches(0.4), Inches(1.4),
                 blurb, size=15, color=CHARCOAL)
        add_bullets(slide, x + Inches(0.3), y + Inches(2.4), col_w - Inches(0.4), Inches(1.8),
                    items, size=13, bullet_color=accent)

    add_footer(slide, n, total)


def slide_toolkit(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Phase 2 — toolkit", "Six draft-only agents around the human.")

    agents = [
        ("Incorporation drafter", "Generates NNC1, articles, first-board resolutions"),
        ("Compliance calendar", "NAR1, BR, BIR51, audit, SCR — tracked and drafted"),
        ("Bookkeeping agent", "Categorises transactions; prepares working papers"),
        ("Profits-tax assistant", "Provisional tax, BIR51 drafts, offshore-source flags"),
        ("Document vault", "Versioned articles, contracts, minutes, filings"),
        ("Decision log", "Every agent action with input, output, approver, time"),
    ]
    col_w = (CONTENT_W - Inches(0.4)) / 2
    for i, (title, body) in enumerate(agents):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (col_w + Inches(0.4))
        y = Inches(2.4) + row * Inches(1.25)
        add_rect(slide, x, y, col_w, Inches(1.1), WHITE)
        add_rect(slide, x, y, Inches(0.12), Inches(1.1), JADE)
        add_text(slide, x + Inches(0.35), y + Inches(0.15), col_w - Inches(0.4), Inches(0.45),
                 title, size=16, color=INK_BLUE, bold=True, font=HEADLINE_FONT)
        add_text(slide, x + Inches(0.35), y + Inches(0.6), col_w - Inches(0.4), Inches(0.5),
                 body, size=13, color=CHARCOAL)

    # guardrail strip
    y = Inches(6.55)
    add_rect(slide, MARGIN, y, CONTENT_W, Inches(0.45), INK_BLUE)
    add_text(slide, MARGIN + Inches(0.3), y + Inches(0.08), CONTENT_W, Inches(0.35),
             "Hard rule: no agent submits to CR, IRD, MPFA, or a bank. Drafts only. Every time.",
             size=12, color=BONE, bold=True, font=MONO_FONT)

    add_footer(slide, n, total)


def slide_moat(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Defensibility", "Where the moat sits.")

    moats = [
        ("HK-anchored curriculum",
         "Cap. 622, IRD, MPFA, two-tier tax — sourced and persona-adapted, not generic."),
        ("Adaptive pedagogy at three personas",
         "Same content, three pacings, three vocabularies, three accessibility defaults."),
        ("Human-in-charge guardrails",
         "Drafts only. Decision log on every agent action. CPA-friendly, not CPA-replacing."),
        ("Compounding artifacts",
         "Each learner produces NNC1, NAR1, audit pack, runbook — reusable, exportable, theirs."),
        ("Distribution wedge",
         "Pre-retirement professionals reached via professional bodies, banks, and family referral."),
    ]
    add_bullets(
        slide, MARGIN, Inches(2.4), CONTENT_W, Inches(4.5),
        [f"{title}  —  {body}" for title, body in moats],
        size=16, spacing=14, bullet_color=AMBER,
    )
    add_footer(slide, n, total)


def slide_business_model(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Business model", "Subscription + per-artifact, never per-filing.")

    tiers = [
        ("Learner",
         "HK$ 88 / mo",
         ["All six arcs", "Adaptive tutor", "Readiness score", "Personal artifacts"],
         JADE),
        ("Operator",
         "HK$ 388 / mo",
         ["Everything in Learner", "Phase 2 toolkit", "Decision log + vault", "Compliance calendar"],
         INK_BLUE),
        ("Operator + CPA",
         "HK$ 988 / mo",
         ["Everything in Operator", "HK-CPA audit-pack review", "Annual filing checklist sign-off", "Priority human escalation"],
         AMBER),
    ]
    col_w = (CONTENT_W - Inches(0.4) * 2) / 3
    for i, (name, price, feats, accent) in enumerate(tiers):
        x = MARGIN + i * (col_w + Inches(0.4))
        y = Inches(2.4)
        add_rect(slide, x, y, col_w, Inches(4.3), WHITE)
        add_rect(slide, x, y, col_w, Inches(0.5), accent)
        add_text(slide, x + Inches(0.3), y + Inches(0.08), col_w, Inches(0.4),
                 name.upper(), size=13, color=BONE, bold=True, font=MONO_FONT)
        add_text(slide, x + Inches(0.3), y + Inches(0.7), col_w - Inches(0.4), Inches(0.7),
                 price, size=28, color=INK_BLUE, bold=True, font=HEADLINE_FONT)
        add_bullets(slide, x + Inches(0.3), y + Inches(1.7), col_w - Inches(0.4), Inches(2.4),
                    feats, size=13, bullet_color=accent, spacing=6)

    add_text(slide, MARGIN, Inches(6.8), CONTENT_W, Inches(0.4),
             "Pricing is never per-filing. We never bill on an action we can't take ourselves.",
             size=11, color=JADE, font=MONO_FONT, bold=True)
    add_footer(slide, n, total)


def slide_market(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Market", "Top-down, HK-anchored.")

    bands = [
        ("TAM", "HK$ 4.8B / yr",
         "All HK professionals 35–75 who could plausibly run a one-person Ltd.",
         INK_BLUE),
        ("SAM", "HK$ 720M / yr",
         "HK professionals with existing client revenue, pre-retirement cohort, and active retirees with marketable expertise.",
         JADE),
        ("SOM (yr 3)", "HK$ 36M / yr",
         "5% of SAM at blended HK$ 300/mo. Conservative on conversion; assumes professional-body and bank distribution partnerships.",
         AMBER),
    ]
    for i, (label, value, body, accent) in enumerate(bands):
        y = Inches(2.4) + i * Inches(1.45)
        add_rect(slide, MARGIN, y, CONTENT_W, Inches(1.3), WHITE)
        add_rect(slide, MARGIN, y, Inches(0.12), Inches(1.3), accent)
        add_text(slide, MARGIN + Inches(0.35), y + Inches(0.2), Inches(1.5), Inches(0.5),
                 label, size=14, color=JADE, bold=True, font=MONO_FONT)
        add_text(slide, MARGIN + Inches(2.0), y + Inches(0.15), Inches(3.0), Inches(0.7),
                 value, size=26, color=INK_BLUE, bold=True, font=HEADLINE_FONT)
        add_text(slide, MARGIN + Inches(5.2), y + Inches(0.2), CONTENT_W - Inches(5.4), Inches(1.0),
                 body, size=13, color=CHARCOAL)

    add_footer(slide, n, total)


def slide_brand(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Brand", "What you'll see, hear, and feel.")

    add_text(slide, MARGIN, Inches(2.4), Inches(6), Inches(0.5),
             "Voice", size=14, color=JADE, bold=True, font=MONO_FONT)
    add_text(slide, MARGIN, Inches(2.85), Inches(6.5), Inches(3.0),
             "Precise, warm, direct.\nLike a respected, patient teacher.\nNot a chatbot mascot.\nNever streaks, never shame, never autofiling.",
             size=18, color=CHARCOAL)

    add_text(slide, MARGIN, Inches(5.6), Inches(6), Inches(0.5),
             "Promise", size=14, color=JADE, bold=True, font=MONO_FONT)
    add_text(slide, MARGIN, Inches(6.05), Inches(6.5), Inches(1.0),
             "From learner to operator.", size=24, color=INK_BLUE, bold=True, font=HEADLINE_FONT)

    # palette swatches
    add_text(slide, Inches(7.5), Inches(2.4), Inches(5), Inches(0.5),
             "Palette", size=14, color=JADE, bold=True, font=MONO_FONT)
    swatches = [
        ("Ink Blue", INK_BLUE, "#0E2A47"),
        ("Jade", JADE, "#2E7D6B"),
        ("Amber", AMBER, "#D89B3A"),
        ("Bone", BONE, "#F6F1E7"),
        ("Charcoal", CHARCOAL, "#1A1A1A"),
    ]
    for i, (name, color, hex_) in enumerate(swatches):
        y = Inches(2.9) + i * Inches(0.72)
        sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), y, Inches(0.6), Inches(0.6))
        set_fill(sw, color)
        add_text(slide, Inches(8.3), y + Inches(0.05), Inches(2.0), Inches(0.3),
                 name, size=14, color=INK_BLUE, bold=True)
        add_text(slide, Inches(8.3), y + Inches(0.3), Inches(2.0), Inches(0.3),
                 hex_, size=11, color=CHARCOAL, font=MONO_FONT)

    add_footer(slide, n, total)


def slide_traction(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Traction & roadmap", "Where we are. Where the next dollar goes.")

    milestones = [
        ("Now",
         "Phase 1 — Arc 1 shipped",
         "Adaptive tutor prototype runs end-to-end. Persona classifier, onboarding dialogue, six-arc spec, canonical Cap. 622 reference, decision-log scaffolding."),
        ("Q+1",
         "Arcs 2–3 + closed beta",
         "Incorporating in HK and Running It. 100 hand-picked learners across the three personas. NPS + readiness-score telemetry."),
        ("Q+2",
         "Arcs 4–6 + toolkit alpha",
         "Agent stack, risk/insurance, succession. Compliance calendar and bookkeeping agent in alpha with 20 paid operators."),
        ("Q+3",
         "Toolkit GA + CPA partner network",
         "Audit-pack review with two HK CPA firms. Bank partnership for KYC narrative."),
    ]
    for i, (when, what, body) in enumerate(milestones):
        y = Inches(2.4) + i * Inches(1.05)
        add_rect(slide, MARGIN, y, CONTENT_W, Inches(0.9), WHITE)
        add_rect(slide, MARGIN, y, Inches(1.0), Inches(0.9), INK_BLUE)
        add_text(slide, MARGIN, y + Inches(0.28), Inches(1.0), Inches(0.4),
                 when, size=13, color=BONE, bold=True, font=MONO_FONT,
                 align=PP_ALIGN.CENTER)
        add_text(slide, MARGIN + Inches(1.2), y + Inches(0.12), Inches(3.6), Inches(0.5),
                 what, size=16, color=INK_BLUE, bold=True, font=HEADLINE_FONT)
        add_text(slide, MARGIN + Inches(1.2), y + Inches(0.45), CONTENT_W - Inches(1.4), Inches(0.5),
                 body, size=12, color=CHARCOAL)

    add_footer(slide, n, total)


def slide_team_ask(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BONE)
    add_section_header(slide, "Team & ask", "What we are raising, and why now.")

    add_text(slide, MARGIN, Inches(2.4), Inches(6.2), Inches(0.5),
             "Team", size=14, color=JADE, bold=True, font=MONO_FONT)
    add_bullets(
        slide, MARGIN, Inches(2.9), Inches(6.2), Inches(3.5),
        [
            "Founder — product lead. HK-based. Built the Phase-1 prototype end-to-end.",
            "Advisory — HK CPA, HK company secretary, HK pension specialist.",
            "Hiring with the raise: senior content lead (HK), staff engineer (agents), designer (accessibility).",
        ],
        size=14, bullet_color=JADE, spacing=10,
    )

    add_text(slide, Inches(7.2), Inches(2.4), Inches(5.5), Inches(0.5),
             "The ask", size=14, color=AMBER, bold=True, font=MONO_FONT)
    add_text(slide, Inches(7.2), Inches(2.9), Inches(5.5), Inches(1.0),
             "USD 1.5M seed", size=32, color=INK_BLUE, bold=True, font=HEADLINE_FONT)
    add_bullets(
        slide, Inches(7.2), Inches(4.0), Inches(5.5), Inches(3.0),
        [
            "18 months runway",
            "Ship Arcs 2–6 with closed-beta cohorts",
            "Toolkit alpha → GA with 2 HK CPA partners",
            "1 bank distribution pilot in HK",
        ],
        size=14, bullet_color=AMBER, spacing=10,
    )
    add_footer(slide, n, total)


def slide_closing(prs, n=None, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, INK_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), AMBER)

    add_text(slide, MARGIN, Inches(2.6), CONTENT_W, Inches(2.4),
             "From learner\nto operator.",
             size=72, color=BONE, bold=True, font=HEADLINE_FONT)
    add_text(slide, MARGIN, Inches(5.4), CONTENT_W, Inches(0.6),
             "OPC Tutor — adaptive AI tutor and agent toolkit for the Hong Kong one-person Ltd.",
             size=18, color=BONE)
    add_text(slide, MARGIN, Inches(6.6), CONTENT_W, Inches(0.4),
             "opc.tutor  ·  Hong Kong  ·  Cap. 622",
             size=12, color=AMBER, font=MONO_FONT)


# ---------- assembly ----------


SLIDES = [
    slide_title,
    slide_pitch,
    slide_two_waves,
    slide_hk_numbers,
    slide_unit,
    slide_what_is_opc,
    slide_product,
    slide_personas,
    slide_toolkit,
    slide_moat,
    slide_brand,
    slide_business_model,
    slide_market,
    slide_traction,
    slide_team_ask,
    slide_closing,
]


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(SLIDES)
    for i, builder in enumerate(SLIDES, start=1):
        builder(prs, i, total)

    out = Path(__file__).parent / "pitch-deck.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}")