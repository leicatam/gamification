#!/usr/bin/env python3
"""Builds the 1.5-hour DTL training deck on valuing deep-tech ventures.

Visual edition: refined typography, infographics for every approach / process /
method, and worked numeric examples beside every formula.  Body and supporting
text sized for projection (14 pt+ for prose, 12-13 pt for dense tables).

Run:     python3 build_deck.py
Output:  Valuing_Deep_Tech_Ventures_DTL_Training.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ----------------------------------------------------------------------
# Theme & typography
# ----------------------------------------------------------------------
NAVY   = RGBColor(0x0F, 0x2A, 0x4A)
NAVY2  = RGBColor(0x16, 0x3A, 0x5E)
NAVY3  = RGBColor(0x21, 0x4A, 0x73)
ACCENT = RGBColor(0xE0, 0x8E, 0x2B)
ACCENT2= RGBColor(0xF2, 0xB5, 0x5C)
TEAL   = RGBColor(0x1B, 0x9E, 0x8A)
GREEN  = RGBColor(0x3C, 0x8C, 0x55)
RED    = RGBColor(0xC0, 0x4A, 0x3A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xEE, 0xF1, 0xF4)
PANEL  = RGBColor(0xEC, 0xF0, 0xF4)
SAND   = RGBColor(0xFB, 0xF1, 0xDF)
GREY   = RGBColor(0x4C, 0x5D, 0x6D)
DARK   = RGBColor(0x1F, 0x2A, 0x36)

HEAD = "Georgia"
BODY = "Calibri"
MONO = "Consolas"

STEP_COLORS = [
    RGBColor(0x4E, 0x79, 0xA8), RGBColor(0x3F, 0x8E, 0x9C),
    RGBColor(0x2E, 0x9E, 0x8A), RGBColor(0x6F, 0xA0, 0x4A),
    RGBColor(0xC9, 0x96, 0x2E), RGBColor(0xE0, 0x8E, 0x2B),
]

DECK_NAME = "Valuing Deep-Tech Ventures  ·  Deep-Tech Lab (DTL) Training"

RECT  = MSO_SHAPE.RECTANGLE
RRECT = MSO_SHAPE.ROUNDED_RECTANGLE
OVAL  = MSO_SHAPE.OVAL
TRI   = MSO_SHAPE.ISOSCELES_TRIANGLE
CHEV  = MSO_SHAPE.CHEVRON
RARR  = MSO_SHAPE.RIGHT_ARROW
DARR  = MSO_SHAPE.DOWN_ARROW
L, R, C = PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.CENTER
MID, TOP = MSO_ANCHOR.MIDDLE, MSO_ANCHOR.TOP

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_counter = {"n": 0}


# ----------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------
def T(text, size=14, color=DARK, bold=False, italic=False, name=BODY,
      align=L, sa=0, sb=0):
    return (text, dict(size=size, color=color, bold=bold, italic=italic,
                       name=name, align=align, sa=sa, sb=sb))


def _fill_tf(tf, lines, anchor=TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw["align"]
        p.space_after = Pt(kw["sa"])
        p.space_before = Pt(kw["sb"])
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(kw["size"]); f.bold = kw["bold"]; f.italic = kw["italic"]
        f.color.rgb = kw["color"]; f.name = kw["name"]


def tbox(slide, l, t, w, h, lines, anchor=TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    _fill_tf(tf, lines, anchor)
    return tb


def shp(slide, kind, l, t, w, h, fill=None, line=None, lw=1.0, rad=None):
    s = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    if rad is not None and kind == RRECT:
        try:
            s.adjustments[0] = rad
        except Exception:
            pass
    return s


def shp_text(s, lines, anchor=MID, ml=0.07, mt=0.04):
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Inches(ml)
    tf.margin_top = tf.margin_bottom = Inches(mt)
    _fill_tf(tf, lines, anchor)
    return s


def connector(slide, x1, y1, x2, y2, color=GREY, w=1.25, dash=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    if dash:
        try:
            from pptx.oxml.ns import qn
            ln = cn.line._get_or_add_ln()
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
        except Exception:
            pass
    return cn


def footer(slide):
    _counter["n"] += 1
    tbox(slide, 0.55, 7.08, 9.5, 0.30, [T(DECK_NAME, 9, GREY)])
    tbox(slide, 12.0, 7.0, 0.8, 0.4, [T(str(_counter["n"]), 12.5, NAVY,
                                        bold=True, name=HEAD, align=R)])


def header(slide, kicker, title):
    shp(slide, RECT, 0, 0, 13.333, 1.06, NAVY)
    shp(slide, RECT, 0, 1.06, 13.333, 0.055, ACCENT)
    tbox(slide, 0.55, 0.13, 12.0, 0.32, [T(kicker.upper(), 12, ACCENT, bold=True)])
    tbox(slide, 0.55, 0.40, 12.25, 0.64, [T(title, 25, WHITE, bold=True, name=HEAD)])


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def callout(slide, top, body, label="KEY IDEA", h=0.92):
    shp(slide, RRECT, 0.55, top, 12.23, h, PANEL, rad=0.10)
    shp(slide, RECT, 0.55, top, 0.11, h, ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.84), Inches(top), Inches(11.8), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MID
    p = tf.paragraphs[0]
    p.alignment = L
    r = p.add_run(); r.text = label + "    "
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = ACCENT
    r.font.name = BODY
    r = p.add_run(); r.text = body
    r.font.size = Pt(14); r.font.color.rgb = DARK; r.font.name = BODY


def bullets(slide, items, top=1.40, height=5.2, lh=16):
    tb = slide.shapes.add_textbox(Inches(0.62), Inches(top), Inches(12.1), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        lvl, text = item[0], item[1]
        bold = item[2] if len(item) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if lvl == 0:
            p.space_after = Pt(10); p.space_before = Pt(2)
            r = p.add_run(); r.text = "▪  "
            r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ACCENT
            r.font.name = BODY
            r = p.add_run(); r.text = text
            r.font.size = Pt(lh); r.font.color.rgb = DARK; r.font.bold = bold
            r.font.name = BODY
        else:
            p.space_after = Pt(6)
            r = p.add_run(); r.text = "        –  "
            r.font.size = Pt(14); r.font.color.rgb = ACCENT; r.font.name = BODY
            r = p.add_run(); r.text = text
            r.font.size = Pt(14); r.font.color.rgb = GREY; r.font.name = BODY


def table(slide, headers, rows, top, col_widths, fsize=13, height=None):
    ncols, nrows = len(headers), len(rows) + 1
    total_w = sum(col_widths)
    left = (13.333 - total_w) / 2
    if height is None:
        height = 0.50 + 0.44 * len(rows)
    gtab = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                  Inches(total_w), Inches(height))
    tbl = gtab.table
    tbl.first_row = True
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MID
        c.margin_left = c.margin_right = Inches(0.10)
        c.margin_top = c.margin_bottom = Inches(0.05)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.size = Pt(fsize + 0.5); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = HEAD
    for i, rowdata in enumerate(rows):
        for j, val in enumerate(rowdata):
            c = tbl.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT
            c.vertical_anchor = MID
            c.margin_left = c.margin_right = Inches(0.10)
            c.margin_top = c.margin_bottom = Inches(0.045)
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fsize); r.font.color.rgb = DARK
            r.font.bold = (j == 0); r.font.name = BODY
    return tbl


def new_slide():
    return prs.slides.add_slide(BLANK)


# ----------------------------------------------------------------------
# Composite slide builders
# ----------------------------------------------------------------------
def content_slide(kicker, title, items, callout_text=None, notes_text="",
                  lh=16):
    s = new_slide()
    header(s, kicker, title)
    if callout_text:
        bullets(s, items, top=1.46, height=3.9, lh=lh)
        callout(s, 5.74, callout_text, h=0.92)
    else:
        bullets(s, items, top=1.56, height=5.2, lh=lh)
    footer(s)
    notes(s, notes_text)
    return s


def table_slide(kicker, title, lead, headers, rows, col_widths,
                callout_text=None, notes_text="", fsize=13, ctop=6.04):
    s = new_slide()
    header(s, kicker, title)
    top = 1.42
    if lead:
        tbox(s, 0.62, 1.24, 12.1, 0.5, [T(lead, 14, GREY, italic=True)])
        top = 1.84
    table(s, headers, rows, top, col_widths, fsize=fsize)
    if callout_text:
        callout(s, ctop, callout_text, h=7.0 - ctop - 0.1)
    footer(s)
    notes(s, notes_text)
    return s


def diagram_slide(kicker, title, notes_text=""):
    s = new_slide()
    header(s, kicker, title)
    footer(s)
    notes(s, notes_text)
    return s


def divider(no, kicker, title, subtitle, notes_text=""):
    s = new_slide()
    shp(s, RECT, 0, 0, 13.333, 7.5, NAVY)
    shp(s, RECT, 0, 0, 13.333, 0.18, ACCENT)
    tbox(s, 0.80, 1.35, 11.0, 2.6, [T(no, 165, NAVY2, bold=True, name=HEAD)])
    shp(s, RECT, 0.95, 4.08, 1.7, 0.12, ACCENT)
    tbox(s, 0.97, 3.56, 11.0, 0.4, [T(kicker.upper(), 15, ACCENT, bold=True)])
    tbox(s, 0.93, 4.28, 11.5, 1.2, [T(title, 37, WHITE, bold=True, name=HEAD)])
    tbox(s, 0.97, 5.44, 11.3, 1.1,
         [T(subtitle, 16.5, RGBColor(0xB5, 0xC7, 0xD8))])
    footer(s)
    notes(s, notes_text)
    return s


def formula_slide(kicker, title, formula_lines, legend, example_lines,
                  callout_text, notes_text=""):
    s = new_slide()
    header(s, kicker, title)
    shp(s, RRECT, 1.05, 1.42, 11.23, 1.12, NAVY2, rad=0.12)
    tf_lines = []
    for i, ln in enumerate(formula_lines):
        tf_lines.append(T(ln, 19 if i == 0 else 14.5,
                          WHITE if i == 0 else ACCENT2,
                          bold=(i == 0), name=MONO, align=C,
                          sa=5 if i == 0 else 1))
    tbox(s, 1.25, 1.42, 10.83, 1.12, tf_lines, anchor=MID)
    # legend (left)
    shp(s, RECT, 0.62, 2.82, 0.08, 0.32, ACCENT)
    tbox(s, 0.82, 2.78, 5.0, 0.4, [T("WHERE", 13, NAVY, bold=True)])
    leg = []
    for sym, desc in legend:
        leg.append(T(sym, 14, ACCENT, bold=True, name=MONO, sa=2))
        leg.append(T(desc, 13, DARK, sa=12))
    tbox(s, 0.82, 3.18, 5.55, 2.7, leg)
    # worked example panel (right)
    shp(s, RRECT, 6.62, 2.72, 6.10, 3.25, SAND, rad=0.05)
    shp(s, RECT, 6.62, 2.72, 6.10, 0.46, ACCENT)
    tbox(s, 6.84, 2.72, 5.7, 0.46,
         [T("WORKED EXAMPLE  —  numbers in the formula", 12.5, NAVY,
            bold=True)], anchor=MID)
    ex = []
    for txt, strong in example_lines:
        ex.append(T(txt, 14 if strong else 13.5,
                    NAVY if strong else DARK, bold=strong,
                    name=MONO if strong else BODY, sa=6))
    tbox(s, 6.90, 3.30, 5.6, 2.55, ex)
    callout(s, 6.06, callout_text, h=0.84)
    footer(s)
    notes(s, notes_text)
    return s


# ----------------------------------------------------------------------
# Infographic primitives
# ----------------------------------------------------------------------
def step_chart(slide, x0, y_base, width, total_h, values, labels, max_val):
    n = len(values)
    sw = width / n
    scale = total_h / max_val
    connector(slide, x0, y_base, x0 + width, y_base, GREY, 1.25)
    for i, v in enumerate(values):
        h = max(v * scale, 0.06)
        x = x0 + i * sw
        col = STEP_COLORS[min(i, len(STEP_COLORS) - 1)]
        shp(slide, RECT, x + 0.06, y_base - h, sw - 0.12, h, col)
        tbox(slide, x, y_base - h - 0.34, sw, 0.30,
             [T(f"${v}M", 13, NAVY, bold=True, name=MONO, align=C)])
        tbox(slide, x, y_base + 0.08, sw, 0.50,
             [T(labels[i], 12, DARK, bold=True, align=C)])


def chain_flow(slide, items, top, h=1.05, gap=0.36):
    n = len(items)
    total = 12.1
    pw = (total - (n - 1) * gap) / n
    x = 0.62
    for i, (title, sub, col) in enumerate(items):
        p = shp(slide, RRECT, x, top, pw, h, col, rad=0.13)
        shp_text(p, [T(title, 14, WHITE, bold=True, name=HEAD, align=C, sa=3),
                     T(sub, 11.5, RGBColor(0xE6, 0xEC, 0xF2), align=C)])
        if i < n - 1:
            shp(slide, CHEV, x + pw + 0.03, top + h / 2 - 0.17,
                gap - 0.05, 0.34, ACCENT)
        x += pw + gap


def ownership_bar(slide, x, y, w, h, segments):
    cx = x
    for label, pct, col in segments:
        seg_w = w * pct
        shp(slide, RECT, cx, y, seg_w, h, col)
        tb = slide.shapes.add_textbox(Inches(cx), Inches(y), Inches(seg_w), Inches(h))
        tf = tb.text_frame
        tf.vertical_anchor = MID
        _fill_tf(tf, [T(label, 12, WHITE, bold=True, align=C, sa=1),
                      T(f"{round(pct*100)}%", 15, WHITE, bold=True,
                        name=HEAD, align=C)], MID)
        cx += seg_w


# ======================================================================
# SLIDE 1 — Title
# ======================================================================
s = new_slide()
shp(s, RECT, 0, 0, 13.333, 7.5, NAVY)
shp(s, RECT, 0, 0, 0.34, 7.5, ACCENT)
for d in [3.4, 2.5, 1.6]:
    shp(s, OVAL, 11.6 - d / 2, 5.0 - d / 2, d, d, None, line=NAVY3, lw=1.4)
tbox(s, 1.1, 1.26, 11.0, 0.45,
     [T("DEEP-TECH LAB (DTL)   ·   TRAINING SERIES", 16, ACCENT, bold=True)])
tbox(s, 1.05, 1.88, 11.0, 1.3,
     [T("Valuing Deep-Tech Ventures", 50, WHITE, bold=True, name=HEAD)])
tbox(s, 1.07, 2.92, 10.9, 0.7,
     [T("A Practitioner's Framework for Researchers & Scientists", 23,
        RGBColor(0xC6, 0xD4, 0xE2))])
shp(s, RECT, 1.12, 3.92, 3.0, 0.10, ACCENT)
tbox(s, 1.1, 4.26, 10.7, 1.6, [
    T("90-minute training  ·  six connected methods, assembled into one "
      "master model.", 15, RGBColor(0xB5, 0xC7, 0xD8), sa=6),
    T("From milestone roadmap and use-of-funds, through benchmarking and "
      "uncertainty, to stakeholder rights and dilution — each illustrated "
      "and worked through with numbers.", 15, RGBColor(0xB5, 0xC7, 0xD8)),
])
tbox(s, 1.1, 6.35, 11.0, 0.5,
     [T("Audience: research scientists and venture founders   |   "
        "All figures are illustrative worked examples, not investment advice.",
        12.5, GREY, italic=True)])
notes(s, "Welcome. This 90-minute session turns valuation from a black box into "
         "a constructed, defensible argument. Every method is shown as a diagram "
         "and worked through with numbers.")
_counter["n"] = 1

# ======================================================================
# SLIDE 2 — The valuation gap
# ======================================================================
content_slide(
    "Why this training", "The valuation gap for scientist-founders",
    [
        (0, "Deep-tech ventures are valued years before revenue — earnings, "
            "book value and cash flow are near zero, so standard accounting "
            "methods return little that is meaningful."),
        (0, "Yet valuation is unavoidable: it sets how much equity you exchange "
            "for the capital needed to reach your next scientific milestone."),
        (0, "A defensible valuation is not a guess — it is an argument that "
            "links three things:"),
        (1, "the science you will do (milestones),"),
        (1, "the funds you will spend to do it (use of funds), and"),
        (1, "the market evidence that prices the result (benchmarking)."),
        (0, "This session builds that argument in six connected methods, each "
            "explained with a diagram and a worked example."),
    ],
    callout_text="A deep-tech valuation is a risk-adjusted, milestone-indexed "
                 "estimate of future value — not a statement of present fact.",
    notes_text="Scientists often feel valuation is arbitrary. The message: it is "
               "a structured argument linking milestones, funds and market "
               "evidence.")

# ======================================================================
# SLIDE 3 — Agenda
# ======================================================================
s = diagram_slide("90-minute run of show", "Agenda — the path of the session")
agenda = [
    ("1", "The milestone roadmap to a valuation", "10 min", STEP_COLORS[0]),
    ("2", "Use of funds → value creation → the math", "20 min", STEP_COLORS[1]),
    ("3", "Benchmarking: market, competition, references", "15 min", STEP_COLORS[2]),
    ("4", "Matching: use of funds tested against the market", "10 min", STEP_COLORS[3]),
    ("5", "Uncertainty: downside / upside & the Golden Triangle", "15 min", STEP_COLORS[4]),
    ("6", "Stakeholders: rights, protections, control, dilution", "12 min", STEP_COLORS[5]),
    ("7", "The master model + comparison of other methods", "8 min", NAVY3),
]
y = 1.36
connector(s, 1.04, 1.66, 1.04, 6.18, RGBColor(0xC9, 0xD3, 0xDC), 2.0)
for num, txt, dur, col in agenda:
    shp(s, OVAL, 0.71, y, 0.62, 0.62, col)
    tbox(s, 0.71, y, 0.62, 0.62,
         [T(num, 19, WHITE, bold=True, name=HEAD, align=C)], anchor=MID)
    card = shp(s, RRECT, 1.62, y, 8.75, 0.62, LIGHT, rad=0.18)
    shp_text(card, [T(txt, 15, DARK, bold=True, name=HEAD)], ml=0.24)
    pill = shp(s, RRECT, 10.60, y + 0.06, 1.68, 0.50, NAVY, rad=0.5)
    shp_text(pill, [T(dur, 12, WHITE, bold=True, align=C)])
    y += 0.70
callout(s, 6.32, "Objectives — by the end you can build a milestone-indexed "
        "valuation, defend it with benchmarking, state its downside and upside, "
        "and read how a term sheet changes your value.", h=0.70)
notes(s, "Walk the arc. Each section builds on the previous one and becomes a "
         "layer of the master model in Section 7.")

# ======================================================================
# SLIDE 4 — Framework at a glance
# ======================================================================
s = diagram_slide("The framework", "The six methods, and how they combine")
methods = [
    ("01", "Milestone roadmap", "when value is created", STEP_COLORS[0]),
    ("02", "Use of funds → rNPV", "bottom-up intrinsic value", STEP_COLORS[1]),
    ("03", "Benchmarking", "top-down market evidence", STEP_COLORS[2]),
    ("04", "Matching method", "reconcile the two views", STEP_COLORS[3]),
    ("05", "Uncertainty", "downside vs upside range", STEP_COLORS[4]),
    ("06", "Stakeholder terms", "value per stakeholder", STEP_COLORS[5]),
]
cw, ch, gx, gy = 3.78, 1.14, 0.30, 0.22
x0 = 0.95
for i, (no, name, sub, col) in enumerate(methods):
    cx = x0 + (i % 3) * (cw + gx)
    cy = 1.30 + (i // 3) * (ch + gy)
    shp(s, RRECT, cx, cy, cw, ch, WHITE, line=col, lw=2.0, rad=0.10)
    shp(s, RECT, cx, cy, 0.64, ch, col)
    tbox(s, cx, cy, 0.64, ch, [T(no, 17, WHITE, bold=True, name=HEAD, align=C)],
         anchor=MID)
    tbox(s, cx + 0.80, cy, cw - 0.94, ch, [
        T(name, 15, NAVY, bold=True, name=HEAD, sa=3),
        T(sub, 12.5, GREY, italic=True),
    ], anchor=MID)
    shp(s, DARR, cx + cw / 2 - 0.10, 3.96, 0.20, 0.22, ACCENT)
mb = shp(s, RRECT, 0.95, 4.24, 11.42, 0.80, NAVY, rad=0.14)
shp_text(mb, [T("THE MASTER MODEL   ·   the disciplined assembly of all six "
                "methods into one argument", 15, WHITE, bold=True, name=HEAD,
                align=C)])
shp(s, DARR, 6.56, 5.10, 0.22, 0.24, ACCENT)
ob = shp(s, RRECT, 1.95, 5.44, 9.4, 0.64, ACCENT, rad=0.5)
shp_text(ob, [T("OUTPUT:  a valuation range  ·  a risk profile  ·  an "
                "ownership outcome", 14, NAVY, bold=True, align=C)])
callout(s, 6.24, "Read this as the map of the next 90 minutes — we build "
        "each block, then assemble it.", label="ORIENTATION", h=0.72)
notes(s, "This is the visual contents page. Refer back to it at each section "
         "divider so the audience always knows where they are.")

# ======================================================================
# SECTION 1
# ======================================================================
divider("01", "Section One", "The Milestone Roadmap to a Valuation",
        "Valuation is a process indexed to scientific and commercial milestones "
        "— not a single number fixed at one date.",
        "Section 1 establishes the mental model: value is created step by step "
        "as risk is removed at each milestone.")

# ---- valuation is a process + TRL ladder ----
s = diagram_slide("Section 1  ·  The milestone roadmap",
                  "Valuation is a process, not a number")
bullets(s, [
    (0, "Map the venture on the Technology Readiness Level (TRL) ladder — "
        "from basic principles (TRL 1) to a system proven in operation (TRL 9)."),
    (0, "Each rung pairs a readiness step with a value inflection: a point at "
        "which a specific risk is removed and the venture is worth more."),
    (0, "The roadmap defines WHEN value is created; later sections quantify "
        "HOW MUCH."),
], top=1.30, height=2.0, lh=15)
y_base = 6.05
bw, gap = 1.12, 0.13
x0 = 0.97
for i in range(9):
    h = 0.42 + i * 0.272
    x = x0 + i * (bw + gap)
    col = STEP_COLORS[0] if i < 3 else (TEAL if i < 6 else ACCENT)
    shp(s, RECT, x, y_base - h, bw, h, col)
    tbox(s, x, y_base - 0.36, bw, 0.30,
         [T(f"TRL {i+1}", 11, WHITE, bold=True, align=C)])
tbox(s, x0, y_base + 0.07, 3.4, 0.3,
     [T("Concept · TRL 1–3", 11.5, STEP_COLORS[0], bold=True, align=C)])
tbox(s, x0 + 3 * (bw + gap), y_base + 0.07, 3.4, 0.3,
     [T("Prototype & validation · TRL 4–6", 11.5, TEAL, bold=True, align=C)])
tbox(s, x0 + 6 * (bw + gap), y_base + 0.07, 3.4, 0.3,
     [T("Demonstration & deployment · TRL 7–9", 11.5, ACCENT, bold=True,
        align=C)])
for idx, lab in [(3, "Prototype"), (5, "Pilot data"), (7, "Clearance")]:
    px = x0 + idx * (bw + gap) + bw / 2
    ph = 0.42 + idx * 0.272
    shp(s, OVAL, px - 0.09, y_base - ph - 0.32, 0.18, 0.18, NAVY)
    tbox(s, px - 0.95, y_base - ph - 0.60, 1.9, 0.26,
         [T(lab, 11, NAVY, bold=True, align=C)])
notes(s, "Scientists already think in TRL terms. The pins mark value-inflection "
         "milestones; the Illuminatio MRI venture targets a 510(k) clearance.")

# ---- roadmap table ----
table_slide(
    "Section 1  ·  The milestone roadmap",
    "From idea to value — the milestone ladder",
    "Illustrative; each venture customises its milestones and valuation method.",
    ["Stage", "TRL", "Defining milestone", "Risk removed", "Primary method"],
    [
        ["Concept", "1–3", "Proof of principle, IP filed", "Scientific basis", "Qualitative scorecard"],
        ["Prototype", "4–6", "Working prototype / validated assay", "Engineering feasibility", "Use-of-funds rNPV"],
        ["Validation", "7–8", "Clinical or field validation data", "Efficacy / performance", "rNPV + comparables"],
        ["Regulatory", "8", "Clearance / approval (e.g. 510(k))", "Regulatory risk", "Benchmarked exit value"],
        ["Commercial", "9", "First revenue / signed contracts", "Market acceptance", "Revenue / earnings multiples"],
        ["Scaling", "9", "Repeatable growth, positive margin", "Execution / scale", "DCF + market multiples"],
    ],
    [1.35, 0.75, 3.4, 2.35, 3.25], fsize=12,
    callout_text="The valuation method follows the venture's TRL stage — "
                 "applying a revenue multiple to a pre-revenue venture is a "
                 "category error.",
    notes_text="Method choice is stage-dependent. We revisit this matrix in "
               "Section 7.")

# ---- de-risking = value creation + step chart ----
s = diagram_slide("Section 1  ·  The milestone roadmap",
                  "De-risking equals value creation")
bullets(s, [
    (0, "Each milestone carries a probability of success (POS). Achieving it "
        "removes that risk — and value jumps by the factor 1 / POS."),
    (0, "Value is created in steps, each coinciding with a milestone — and "
        "the early steps are small."),
], top=1.28, height=1.25, lh=15)
tbox(s, 1.30, 2.60, 5.2, 0.32,
     [T("Risk-adjusted venture value (US$M, illustrative)", 12, GREY,
        italic=True)])
step_chart(s, 1.15, 6.02, 11.0, 3.10,
           [4.4, 7.9, 19.7, 41.0, 64.0, 80.0],
           ["Today", "After M1", "After M2", "After M3", "After M4", "Exit Yr 5"],
           80.0)
note = shp(s, RRECT, 1.34, 3.06, 4.55, 1.40, PANEL, rad=0.10)
shp(s, RECT, 1.34, 3.06, 0.10, 1.40, ACCENT)
shp_text(note, [
    T("READ THE CHART", 12, ACCENT, bold=True, sa=5),
    T("Value is back-loaded: most of it appears only as the final risks are "
      "removed. A round's job is to buy the milestones that remove the most "
      "risk per dollar.", 13.5, DARK)], anchor=MID, ml=0.18)
notes(s, "The chart deliberately shows tiny early steps — that IS the "
         "deep-tech reality. The 4.4 → 80 series recurs through the deck.")

# ======================================================================
# SECTION 2
# ======================================================================
divider("02", "Section Two", "From Use of Funds to Value — the Math",
        "Funds drive tasks, tasks remove risk, removed risk accumulates into "
        "valuation — and the rNPV formula makes it measurable.",
        "Section 2 is the quantitative heart of the training.")

# ---- use of funds step 1 ----
s = diagram_slide("Section 2  ·  Use of funds",
                  "Step 1 — use of funds as the unit of analysis")
bullets(s, [
    (0, "Start not with a valuation number, but with the use of funds: the "
        "itemised plan of what the next round of capital will be spent on."),
    (0, "A credible use of funds is structured around milestones, not line "
        "items — it states which milestones the money will deliver."),
    (0, "Each category funds tasks; each task is chosen to remove a specific, "
        "nameable risk."),
], top=1.32, height=2.1, lh=15.5)
cats = [
    ("R&D &\nprototyping", STEP_COLORS[0]),
    ("Validation &\nclinical / pilot", STEP_COLORS[1]),
    ("Regulatory\n& IP", STEP_COLORS[2]),
    ("Team &\noperations", STEP_COLORS[3]),
    ("Go-to-\nmarket", STEP_COLORS[4]),
]
cwid = 2.18
x = 0.85
for name, col in cats:
    box = shp(s, RRECT, x, 3.62, cwid, 1.20, col, rad=0.12)
    shp_text(box, [T(l, 13.5, WHITE, bold=True, name=HEAD, align=C)
                   for l in name.split("\n")])
    x += cwid + 0.20
tbox(s, 0.85, 4.96, 11.5, 0.34,
     [T("Typical use-of-funds categories — every dollar should be "
        "traceable to a milestone", 13, GREY, italic=True, align=C)])
callout(s, 5.50, "The use of funds is the bridge between the laboratory and the "
        "valuation: it is the list of risks the round will buy down.", h=0.92)
notes(s, "Insist the use of funds is milestone-led. 'We will spend $2.5M to "
         "reach prototype and validation' — not just 'we will spend $2.5M'.")

# ---- how funded tasks create value (causal chain) ----
s = diagram_slide("Section 2  ·  Value creation",
                  "How funded tasks create value")
bullets(s, [
    (0, "Trace the causal chain for every dollar. A task creates value only if "
        "it changes a probability of success or the size of the prize."),
    (0, "Tasks that merely maintain the status quo are necessary cost — they "
        "keep the option alive, but do not by themselves create value."),
], top=1.30, height=1.35, lh=15.5)
chain_flow(s, [
    ("Funds", "US$2.5M round", STEP_COLORS[0]),
    ("Funded task", "clinical study", STEP_COLORS[1]),
    ("Milestone", "validation data", STEP_COLORS[2]),
    ("Risk removed", "efficacy risk", STEP_COLORS[3]),
    ("Value created", "POS 0.3 → 0.5", STEP_COLORS[5]),
], top=2.92, h=1.10)
tbox(s, 0.62, 4.24, 12.1, 0.34,
     [T("The value-creation chain — worked for one funded task", 13, GREY,
        italic=True, align=C)])
callout(s, 4.78, "Test every funded task with one question: which probability "
        "does it move, and by how much? If the answer is ‘none’, it is "
        "overhead, not investment.", h=1.02)
notes(s, "This reframes budgeting for scientists: spending is investment only "
         "when it shifts a POS or expands the market.")

# ---- how value accumulates ----
s = diagram_slide("Section 2  ·  Value accumulation",
                  "How value accumulates across milestones")
tbox(s, 0.62, 1.28, 12.1, 0.4,
     [T("Milestones compound: achieving each one removes its risk and "
        "multiplies value by 1 / POS.", 14, GREY, italic=True)])
blocks = [
    ("Today", "$4.4M", "4 milestones remain|CPOS = 16.8%", STEP_COLORS[0]),
    ("After M1", "$7.9M", "prototype done|CPOS = 24%", STEP_COLORS[2]),
    ("After M2", "$19.7M", "validation done|CPOS = 48%", STEP_COLORS[5]),
]
bw, bh = 3.05, 1.95
x = 0.95
for i, (title, val, sub, col) in enumerate(blocks):
    box = shp(s, RRECT, x, 2.30, bw, bh, col, rad=0.10)
    lines = [T(title, 14, WHITE, bold=True, name=HEAD, align=C, sa=3),
             T(val, 30, WHITE, bold=True, name=HEAD, align=C, sa=4)]
    for sl in sub.split("|"):
        lines.append(T(sl, 12.5, RGBColor(0xE8, 0xEE, 0xF3), align=C, sa=1))
    shp_text(box, lines)
    if i < 2:
        ax = x + bw + 0.04
        shp(s, RARR, ax, 2.30 + bh / 2 - 0.32, 0.80, 0.64, ACCENT)
        note = ["M1 achieved|removes POS 0.70|→ value × 1.8",
                "M2 achieved|removes POS 0.50|→ value × 2.5"][i]
        tbox(s, ax - 0.32, 4.40, 1.50, 1.0,
             [T(n, 11, NAVY, bold=True, align=C, sa=1)
              for n in note.split("|")])
    x += bw + 0.92
callout(s, 5.66, "A US$2.5M round that funds M1 + M2 lifts risk-adjusted value "
        "from ~US$4.4M to ~US$19.7M — about 4.5×. That multiple, not the "
        "cash, is the valuation argument.", h=0.92)
notes(s, "The uplift combines risk removed (CPOS 16.8% → 48%) and time "
         "shortened (5 yrs → 3).")

# ---- rNPV formula + worked example ----
formula_slide(
    "Section 2  ·  The math",
    "The math of valuation — risk-adjusted NPV (rNPV)",
    [
        "rNPV₀  =  Σ ( CPOSₜ × CFₜ ) / (1+r)ₜ   −   Σ Costₜ / (1+r)ₜ",
        "value step on milestone success:   V → V × ( 1 / POS )",
    ],
    [
        ("CPOSₜ", "cumulative probability of reaching year t"),
        ("CFₜ", "cash flow / exit value in year t (from benchmarking)"),
        ("Costₜ", "use of funds deployed in year t"),
        ("r", "venture discount rate (15–40% for early deep tech)"),
    ],
    [
        ("Diagnostic-imaging venture, US$2.5M round", True),
        ("1 · CPOS = 0.70×0.50×0.60×0.80 = 0.168", False),
        ("2 · Exit value  CF₅ = US$80M  (Year 5)", False),
        ("3 · Discount  r = 25%  →  1.25⁵ = 3.05", False),
        ("4 · Reward = 0.168 × 80 ÷ 3.05 = US$4.40M", True),
        ("5 · Cost term = PV of US$2.5M ≈ US$2.15M", False),
        ("rNPV₀ = 4.40 − 2.15 ≈ US$2.25M", True),
        ("    ( > 0  →  the round creates value )", False),
    ],
    "rNPV (also called eNPV) is the standard intrinsic method for risk-laden, "
    "milestone-driven ventures — it makes de-risking measurable.",
    notes_text="First sum = risk-adjusted reward; second = cost. Discipline: "
               "once a milestone is achieved, DROP its POS from CPOS — do not "
               "also keep an inflated r, or you double-count the risk.")

# ---- worked example table ----
s = diagram_slide("Section 2  ·  Worked example",
                  "Worked example — a milestone-funded venture")
tbox(s, 0.62, 1.24, 12.1, 0.46,
     [T("Illustrative figures for a diagnostic-imaging venture raising a "
        "US$2.5M round to fund the first two milestones.", 14, GREY,
        italic=True)])
table(s, ["Item", "Value", "Basis"], [
    ["Benchmarked exit value (Year 5)", "US$80M", "comparable med-tech exits"],
    ["M1 — working prototype", "POS 70%", "engineering / feasibility risk"],
    ["M2 — clinical validation data", "POS 50%", "efficacy / scientific risk"],
    ["M3 — regulatory clearance (510(k))", "POS 60%", "regulatory risk"],
    ["M4 — first commercial revenue", "POS 80%", "market-acceptance risk"],
    ["Cumulative POS (all four)", "16.8%", "0.70 × 0.50 × 0.60 × 0.80"],
    ["Risk-adjusted value today (pre-money)", "≈ US$4.4M", "0.168 × 80 ÷ 1.25⁵"],
    ["Value after M1 + M2 achieved", "≈ US$19.7M", "0.48 × 80 ÷ 1.25³"],
], 1.80, [4.6, 1.8, 4.6], fsize=12.5)
callout(s, 6.02, "A US$2.5M round that buys M1 + M2 lifts risk-adjusted value "
        "~4.5× — that multiple, not the cash, is the valuation argument.",
        h=0.88)
notes(s, "Walk each row. 1.25⁵ = 3.05; 1.25³ = 1.95. All figures illustrative.")

# ======================================================================
# SECTION 3
# ======================================================================
divider("03", "Section Three", "The Benchmarking Approach",
        "Anchoring the venture to external evidence: market size, competitive "
        "position, qualitative quality, and current valuation references.",
        "Section 3 supplies the exit value and the sense-checks the rNPV needs.")

# ---- benchmarking 4-lens hub ----
s = diagram_slide("Section 3  ·  Benchmarking",
                  "Benchmarking — triangulating value from the market")
tbox(s, 0.62, 1.26, 12.1, 0.4,
     [T("Four lenses, used together — agreement across all four produces a "
        "defensible exit value.", 14, GREY, italic=True)])
hub_cx, hub_cy, hub_d = 6.665, 4.05, 2.05
lens = [
    (1.05, 1.84, "THE MARKET", "How large is the addressable opportunity?", STEP_COLORS[0]),
    (7.92, 1.84, "COMPETITIVE ANALYSIS", "Who else is there — how do you compare?", STEP_COLORS[2]),
    (1.05, 4.58, "QUALITATIVE ‘GOOD VIEWS’", "Team, IP and execution quality", STEP_COLORS[3]),
    (7.92, 4.58, "VALUATION REFERENCES", "What were comparable ventures valued at?", STEP_COLORS[5]),
]
for lx, ly, name, sub, col in lens:
    connector(s, lx + 4.36 / 2, ly + 0.80, hub_cx, hub_cy,
              RGBColor(0xB7, 0xC3, 0xCE), 1.5)
for lx, ly, name, sub, col in lens:
    shp(s, RRECT, lx, ly, 4.36, 1.60, WHITE, line=col, lw=2.25, rad=0.09)
    shp(s, RECT, lx, ly, 4.36, 0.46, col)
    tbox(s, lx, ly, 4.36, 0.46, [T(name, 13, WHITE, bold=True, align=C)],
         anchor=MID)
    tbox(s, lx + 0.2, ly + 0.56, 3.96, 0.98,
         [T(sub, 13.5, DARK, align=C)], anchor=MID)
hub = shp(s, OVAL, hub_cx - hub_d / 2, hub_cy - hub_d / 2, hub_d, hub_d, NAVY)
shp_text(hub, [T("Defensible", 14, WHITE, bold=True, name=HEAD, align=C, sa=2),
               T("exit value", 14, WHITE, bold=True, name=HEAD, align=C, sa=3),
               T("CFₜ for the rNPV", 11.5, ACCENT2, align=C)])
callout(s, 6.30, "Benchmarking is top-down (market inward); the use-of-funds "
        "rNPV is bottom-up. Section 4 reconciles them.", label="HOW IT FITS",
        h=0.70)
notes(s, "Any single lens can mislead; triangulation across all four makes the "
         "exit value defensible.")

# ---- TAM/SAM/SOM ----
s = diagram_slide("Section 3  ·  The market", "Market sizing — TAM, SAM, SOM")
ccx, ccy = 3.45, 3.95
rings = [(4.05, RGBColor(0xD7, 0xE3, 0xEE), "TAM"),
         (2.78, RGBColor(0x9F, 0xC2, 0xCB), "SAM"),
         (1.52, NAVY, "SOM")]
for d, col, lab in rings:
    shp(s, OVAL, ccx - d / 2, ccy - d / 2, d, d, col)
for d, col, lab in rings:
    tc = WHITE if lab == "SOM" else NAVY
    ty = ccy - d / 2 + (0.13 if lab != "SOM" else d / 2 - 0.17)
    tbox(s, ccx - 0.8, ty, 1.6, 0.3,
         [T(lab, 13 if lab != "SOM" else 15, tc, bold=True, name=HEAD, align=C)])
expl = [
    (1.78, "TAM — Total Addressable Market", STEP_COLORS[0],
     "Total demand if you served everyone.  DTL reference: liver-disease "
     "diagnostics ≈ US$30.6B (2020), ~6.7% CAGR."),
    (3.26, "SAM — Serviceable Available Market", TEAL,
     "The slice your product and channel can reach — e.g. non-invasive "
     "MRI-based liver-fibrosis testing in target markets."),
    (4.74, "SOM — Serviceable Obtainable Market", ACCENT,
     "The share you can realistically win — early hospital and pharma "
     "adopters reachable in the first 3–5 years."),
]
for ey, name, col, body in expl:
    shp(s, RECT, 6.28, ey, 0.09, 1.22, col)
    tbox(s, 6.54, ey, 6.2, 1.26, [
        T(name, 14.5, NAVY, bold=True, name=HEAD, sa=4),
        T(body, 13.5, DARK),
    ])
callout(s, 6.12, "Investors discount an unsupported TAM heavily — a credible "
        "SOM, built bottom-up from named customers, carries far more weight.",
        h=0.74)
notes(s, "Teach moving from the TAM headline to a defensible bottom-up SOM.")

# ---- competitive scorecard ----
table_slide(
    "Section 3  ·  Competition & quality",
    "Competitive analysis and the qualitative ‘good views’",
    "A benchmarking scorecard — the dimensions DTL reviewers use to grade ventures.",
    ["Dimension", "What it benchmarks", "Strong signal"],
    [
        ["Technology content", "Depth, defensibility, IP position", "Granted patents; clear performance edge"],
        ["Team / founder readiness", "Capability and commitment to execute", "Domain track record; full-time founders"],
        ["Commercialisation potential", "Route to revenue and scale", "Named customers; validated demand"],
        ["Feasibility & sustainability", "Can it be delivered and funded", "Credible milestones; realistic capital plan"],
        ["Competitive position", "Advantage vs named rivals", "Direct measurement vs indirect proxies"],
    ],
    [3.1, 4.0, 4.5], fsize=13,
    callout_text="The ‘good views’ are qualitative but not arbitrary — a "
                 "structured scorecard turns judgement into a comparable "
                 "benchmark.",
    notes_text="These five dimensions come from the DTL interview scorecard.")

# ---- valuation references ----
content_slide(
    "Section 3  ·  Valuation references",
    "Current valuation references — reading the comparables",
    [
        (0, "Valuation references are observed prices: recent funding rounds, "
            "acquisitions and licensing deals for comparable ventures."),
        (0, "Three reference types, in decreasing order of availability:"),
        (1, "comparable financing rounds — pre-money valuations at the same stage;"),
        (1, "comparable transactions — acquisition prices for similar assets;"),
        (1, "comparable public companies — listed peers’ trading multiples."),
        (0, "Adjust every reference for stage, geography, TRL and deal terms "
            "before applying it — raw comparables are rarely like-for-like."),
        (0, "Date every reference: venture valuations move with the financing "
            "cycle, so a comparable over 12–18 months old needs caution."),
    ],
    callout_text="A comparable is evidence, not a verdict. Adjust it for stage, "
                 "terms and date — then let it inform a range, not a point.",
    notes_text="Deep-tech comparables are scarce and deal terms seldom "
               "disclosed — hence triangulation.")

# ======================================================================
# SECTION 4
# ======================================================================
divider("04", "Section Four", "The Matching Method",
        "Reconciling the bottom-up use of funds with the top-down benchmarking "
        "— and testing the funding plan against market evidence.",
        "Section 4 is the reconciliation step.")

# ---- matching converge ----
s = diagram_slide("Section 4  ·  Matching",
                  "Matching — reconciling bottom-up and top-down")
tbox(s, 0.62, 1.26, 12.1, 0.4,
     [T("Two independent estimates of value are placed side by side — do "
        "they tell a consistent story?", 14, GREY, italic=True)])
lb = shp(s, RRECT, 0.85, 2.02, 3.55, 1.72, STEP_COLORS[1], rad=0.10)
shp_text(lb, [T("BOTTOM-UP", 12, ACCENT2, bold=True, align=C, sa=2),
              T("Use-of-funds rNPV", 14.5, WHITE, bold=True, name=HEAD,
                align=C, sa=3),
              T("built from milestones, POS|and the spending plan", 12,
                RGBColor(0xE3, 0xE9, 0xEF), align=C)][:2] +
             [T("built from milestones, POS", 12, RGBColor(0xE3,0xE9,0xEF),
                align=C, sa=1),
              T("and the spending plan", 12, RGBColor(0xE3,0xE9,0xEF),
                align=C)])
rb = shp(s, RRECT, 8.93, 2.02, 3.55, 1.72, STEP_COLORS[4], rad=0.10)
shp_text(rb, [T("TOP-DOWN", 12, NAVY, bold=True, align=C, sa=2),
              T("Benchmarked exit value", 14.5, WHITE, bold=True, name=HEAD,
                align=C, sa=3),
              T("built from market size", 12, RGBColor(0xFB,0xF1,0xDF),
                align=C, sa=1),
              T("and comparable references", 12, RGBColor(0xFB,0xF1,0xDF),
                align=C)])
shp(s, RARR, 4.55, 2.60, 1.62, 0.56, STEP_COLORS[1])
shp(s, MSO_SHAPE.LEFT_ARROW, 7.15, 2.60, 1.62, 0.56, STEP_COLORS[4])
hub = shp(s, OVAL, 5.49, 2.06, 2.34, 1.64, NAVY)
shp_text(hub, [T("MATCH", 12, ACCENT, bold=True, align=C, sa=1),
               T("Defensible", 13.5, WHITE, bold=True, name=HEAD, align=C),
               T("valuation range", 13.5, WHITE, bold=True, name=HEAD, align=C)])
oc = shp(s, RRECT, 0.85, 4.06, 5.78, 1.55, RGBColor(0xE4, 0xF0, 0xE7),
         line=GREEN, lw=1.75, rad=0.09)
shp_text(oc, [T("CONVERGE", 13, GREEN, bold=True, align=C, sa=4),
              T("Confidence is high. The overlap defines the defensible "
                "valuation range.", 13, DARK, align=C)])
dv = shp(s, RRECT, 6.75, 4.06, 5.78, 1.55, RGBColor(0xF6, 0xE6, 0xE3),
         line=RED, lw=1.75, rad=0.09)
shp_text(dv, [T("DIVERGE", 13, RED, bold=True, align=C, sa=4),
              T("One assumption is wrong — an optimistic POS, an inflated "
                "TAM, a stale comparable. Find it.", 13, DARK, align=C)])
callout(s, 5.78, "Matching is a consistency test: convergence builds "
        "confidence, divergence locates the flawed assumption.", h=0.90)
notes(s, "The point of matching is to diagnose. A gap tells you where to dig.")

# ---- matching worksheet ----
table_slide(
    "Section 4  ·  Matching", "The matching worksheet",
    "Place each funded milestone against the market evidence that prices it.",
    ["Use of funds (milestone)", "Value mechanism", "Market evidence that tests it", "Match?"],
    [
        ["Prototype build", "Engineering POS 0.4→0.7", "Comparable seed rounds post-prototype", "Aligned"],
        ["Clinical validation study", "Efficacy POS 0.3→0.5", "Valuations after comparable validation", "Aligned"],
        ["Regulatory submission", "Clearance POS → 0.6", "Re-rating of cleared comparable devices", "Aligned"],
        ["Go-to-market / first contract", "Unlocks revenue", "Revenue multiples of early-commercial peers", "Review"],
    ],
    [3.4, 2.7, 4.5, 1.15], fsize=12,
    callout_text="Each row links a dollar of spend to a milestone, to a value "
                 "step, to the market evidence that confirms or challenges it.",
    notes_text="A 'Review' flag means bottom-up and top-down disagree.")

# ---- articulating use of funds ----
content_slide(
    "Section 4  ·  Matching",
    "Articulating use of funds and testing it against the market",
    [
        (0, "Articulate — state the use of funds so each item is explicitly "
            "tied to a milestone, a probability shift and a value step."),
        (0, "Examine each articulated item against the market evidence:"),
        (1, "Is it a milestone the market actually re-prices? (some are invisible);"),
        (1, "Is the probability shift consistent with base rates from comparable programmes?"),
        (1, "Does the implied value step match observed valuations at that stage?"),
        (0, "Where the answers align, the use of funds is validated — the "
            "round is priced by the market, not by assertion."),
        (0, "Where they do not, reallocate funds toward the milestones the "
            "market genuinely rewards."),
    ],
    callout_text="A use of funds is well-articulated when every dollar is "
                 "traceable to a milestone the market is observed to pay for.",
    notes_text="Closes the loop between Sections 2 and 3.")

# ======================================================================
# SECTION 5
# ======================================================================
divider("05", "Section Five", "Uncertainty — Downside vs Upside",
        "A deep-tech valuation is a distribution, not a point. Frame its spread "
        "with the Golden Triangle: technology, organisation and commercial risk.",
        "Section 5 forces the deck away from single numbers.")

# ---- Golden Triangle ----
s = diagram_slide("Section 5  ·  Uncertainty",
                  "The Golden Triangle — the three sources of risk")
tbox(s, 0.62, 1.22, 12.1, 0.36,
     [T("Every deep-tech venture's uncertainty resolves into three vertices "
        "— value is capped by the weakest.", 14, GREY, italic=True)])
apex = (6.675, 2.42)
bl, br = (4.55, 5.12), (8.80, 5.12)
connector(s, 3.85, 2.20, apex[0] - 0.55, apex[1], RGBColor(0xB7, 0xC3, 0xCE), 1.4)
connector(s, 4.10, 5.46, bl[0] - 0.42, bl[1] + 0.28, RGBColor(0xB7, 0xC3, 0xCE), 1.4)
connector(s, 9.22, 5.46, br[0] + 0.42, br[1] + 0.28, RGBColor(0xB7, 0xC3, 0xCE), 1.4)
shp(s, TRI, 3.95, 1.78, 5.45, 3.78, RGBColor(0xE9, 0xEE, 0xF2), line=NAVY, lw=1.5)
ctr = shp(s, OVAL, 5.58, 3.18, 2.18, 1.42, NAVY)
shp_text(ctr, [T("VENTURE", 11, ACCENT, bold=True, align=C, sa=1),
               T("VALUE", 16, WHITE, bold=True, name=HEAD, align=C, sa=2),
               T("capped by the weakest vertex", 10, RGBColor(0xC6, 0xD4, 0xE2),
                 align=C)])
for (vx, vy), short, col in [(apex, "TECH", STEP_COLORS[0]),
                             (bl, "ORG", STEP_COLORS[3]),
                             (br, "COMM", STEP_COLORS[5])]:
    o = shp(s, OVAL, vx - 0.58, vy - 0.58, 1.16, 1.16, col, line=WHITE, lw=2.0)
    shp_text(o, [T(short, 13, WHITE, bold=True, name=HEAD, align=C)])
caps = [
    (0.62, 1.62, 3.18, "TECHNOLOGY", STEP_COLORS[0],
     "Does the science work — at scale, reproducibly, with freedom to operate?"),
    (0.62, 4.86, 3.48, "ORGANISATION", STEP_COLORS[3],
     "Can the team, governance and capital plan execute and survive?"),
    (9.26, 4.86, 3.48, "COMMERCIAL", STEP_COLORS[5],
     "Will customers adopt and pay — against competition and regulation?"),
]
for cx, cy, cwid, name, col, body in caps:
    shp(s, RRECT, cx, cy, cwid, 1.46, WHITE, line=col, lw=2.0, rad=0.10)
    shp(s, RECT, cx, cy, cwid, 0.40, col)
    tbox(s, cx, cy, cwid, 0.40, [T(name, 13, WHITE, bold=True, align=C)],
         anchor=MID)
    tbox(s, cx + 0.16, cy + 0.44, cwid - 0.32, 0.98,
         [T(body, 12.5, DARK)], anchor=MID)
callout(s, 6.44, "Every milestone and risk in the rNPV maps onto one of these "
        "three vertices.", label="WHY IT MATTERS", h=0.58)
notes(s, "Researchers over-weight the technology vertex. Investors price all "
         "three; organisation and commercial risk are where scientist-led "
         "ventures most often fail.")

# ---- scenario fan ----
s = diagram_slide("Section 5  ·  Uncertainty",
                  "Scenario analysis — downside, base and upside")
tbox(s, 0.62, 1.26, 12.1, 0.4,
     [T("Replace the single number with three explicit, fully-described "
        "scenarios.", 14, GREY, italic=True)])
scen = [
    ("UPSIDE", 70, 0.20, GREEN,
     "Faster validation, larger market capture, strategic interest"),
    ("BASE", 25, 0.50, STEP_COLORS[0],
     "Plan delivered roughly on time and on budget"),
    ("DOWNSIDE", 8, 0.30, RED,
     "Key milestone slips; partial recovery via asset sale or pivot"),
]
ox, scale = 2.30, 0.118
y = 1.92
for name, val, prob, col, desc in scen:
    bwid = val * scale
    shp(s, RECT, ox, y, bwid, 0.74, col)
    tbox(s, ox - 1.65, y, 1.50, 0.74,
         [T(name, 13, col, bold=True, align=R)], anchor=MID)
    tbox(s, ox + bwid + 0.14, y, 1.9, 0.74,
         [T(f"US${val}M", 14.5, NAVY, bold=True, name=MONO, sa=1),
          T(f"p = {prob:.2f}", 12, GREY)], anchor=MID)
    tbox(s, ox + 0.14, y + 0.78, 9.0, 0.32,
         [T(desc, 12, GREY, italic=True)])
    y += 1.18
ev_x = ox + 28.9 * scale
connector(s, ev_x, 1.74, ev_x, 5.58, ACCENT, 2.0, dash=True)
tbox(s, ev_x - 1.2, 5.58, 2.4, 0.3,
     [T("E[V] ≈ US$28.9M", 12.5, ACCENT, bold=True, align=C)])
callout(s, 6.04, "Expected value  E[V] = 0.30×8 + 0.50×25 + 0.20×70 ≈ "
        "US$28.9M — the upside tail lifts the mean above the base case.",
        h=0.86)
notes(s, "The deliberate asymmetry — upside far above base — is "
         "characteristic of deep tech.")

# ---- mapping uncertainty across triangle ----
table_slide(
    "Section 5  ·  Uncertainty",
    "Mapping uncertainty across the Golden Triangle",
    "For each vertex, name the specific downside and upside — do not leave risk abstract.",
    ["Vertex", "Downside driver", "Upside driver"],
    [
        ["Technology", "Fails to reproduce or scale; blocked by third-party IP",
         "Performance beats target; platform spawns several products"],
        ["Organisation", "Key person leaves; cash runs out before next milestone",
         "Senior hires land; non-dilutive grants extend the runway"],
        ["Commercial", "Slow adoption; reimbursement tightens; new entrant",
         "Strategic partner or anchor customer; faster path to scale"],
    ],
    [2.4, 5.3, 5.3], fsize=13,
    callout_text="Quantify, do not just list: attach a probability and a value "
                 "impact to each driver so it feeds the scenario weights.",
    notes_text="Each driver should trace to a milestone POS in the rNPV.")

# ---- expected value formula + worked example ----
formula_slide(
    "Section 5  ·  Uncertainty",
    "Expected value and probability weighting",
    [
        "E[V]  =  p_down × V_down  +  p_base × V_base  +  p_up × V_up",
    ],
    [
        ("p_x", "probability of scenario x (from the Golden-Triangle drivers)"),
        ("V_x", "venture value in scenario x — each scenario fully narrated"),
        ("E[V]", "the probability-weighted expected value"),
    ],
    [
        ("Three scenarios for the same venture:", True),
        ("Downside   p = 0.30   V = US$8M", False),
        ("    →  0.30 × 8  =  2.40", False),
        ("Base         p = 0.50   V = US$25M", False),
        ("    →  0.50 × 25  =  12.50", False),
        ("Upside       p = 0.20   V = US$70M", False),
        ("    →  0.20 × 70  =  14.00", False),
        ("E[V] = 2.40 + 12.50 + 14.00 = US$28.9M", True),
    ],
    "Report the range AND the probabilities, not only E[V] — the spread "
    "between downside and upside is itself decision-useful.",
    notes_text="E[V] collapses the distribution to one number — always present "
               "it with the range.")

# ======================================================================
# SECTION 6
# ======================================================================
divider("06", "Section Six",
        "Stakeholders — Rights, Protections, Control & Dilution",
        "Enterprise value is not value received. Term-sheet rights, controls "
        "and future dilution decide who actually gets what.",
        "Section 6 separates the headline valuation from value realised.")

# ---- cap table + ownership bar ----
s = diagram_slide("Section 6  ·  Stakeholders",
                  "The cap table — who holds the value")
bullets(s, [
    (0, "The capitalisation table records who owns the venture — founders, "
        "the employee option pool, and investors by round."),
    (0, "The headline valuation prices the whole company; the cap table and "
        "term sheet decide how that value is divided."),
    (0, "Two distinct dimensions are negotiated: economic rights (who receives "
        "cash, and in what order) and control rights (who can decide or veto)."),
], top=1.34, height=2.3, lh=15.5)
tbox(s, 0.85, 3.74, 11.5, 0.32,
     [T("Illustrative cap table after a seed round", 13, NAVY, bold=True)])
ownership_bar(s, 0.85, 4.08, 11.6, 0.94, [
    ("Founders", 0.55, STEP_COLORS[0]),
    ("Option pool", 0.09, TEAL),
    ("Seed investors", 0.36, ACCENT),
])
callout(s, 5.34, "The headline valuation sizes the pie. Rights, controls and "
        "dilution decide the slice each stakeholder actually eats — "
        "ownership % alone is an incomplete picture.", h=0.92)
notes(s, "A minority investor with strong protective provisions can hold "
         "decisive control.")

# ---- investor rights table ----
table_slide(
    "Section 6  ·  Stakeholders",
    "Investor rights, protections and controls",
    "Common term-sheet provisions — and how each changes the value you receive.",
    ["Provision", "What it does", "Effect on founder value"],
    [
        ["1× liquidation preference", "Investor recovers capital first on exit", "Reduces common value at low/mid exits"],
        ["Participating preferred", "Preference AND a share of the remainder", "‘Double dip’ — cuts founder proceeds"],
        ["Anti-dilution (weighted-avg.)", "Re-prices investor shares in a down round", "Shifts down-round dilution to founders"],
        ["Pro-rata / pre-emptive rights", "Investor may keep its % in later rounds", "Constrains who you can admit later"],
        ["Board seats & protective provisions", "Veto over budget, hiring, financing, sale", "Control independent of ownership %"],
        ["Founder vesting & drag-along", "Re-earn your shares; can be forced to sell", "Aligns — and limits — founder economics"],
    ],
    [3.3, 4.25, 4.45], fsize=12,
    callout_text="Two ventures can carry an identical headline valuation yet "
                 "deliver very different value to founders — the terms decide.",
    notes_text="Participating preferred surprises founders most; 1× "
               "non-participating is the founder-friendly norm.")

# ---- future dilution ----
s = diagram_slide("Section 6  ·  Stakeholders",
                  "Future development and dilution")
tbox(s, 0.62, 1.26, 12.1, 0.4,
     [T("Deep-tech ventures need several rounds — dilution is a planned "
        "sequence, not a one-off event.", 14, GREY, italic=True)])
stages = [
    ("Before seed round", [("Founders", 0.91, STEP_COLORS[0]),
                           ("Pool", 0.09, TEAL)]),
    ("After seed round", [("Founders", 0.55, STEP_COLORS[0]),
                          ("Pool", 0.09, TEAL),
                          ("Seed", 0.36, ACCENT)]),
    ("After Series A", [("Founders", 0.38, STEP_COLORS[0]),
                        ("Pool", 0.10, TEAL),
                        ("Seed", 0.25, ACCENT),
                        ("Series A", 0.27, STEP_COLORS[4])]),
]
y = 1.86
for label, segs in stages:
    tbox(s, 0.85, y - 0.06, 11.6, 0.30, [T(label, 13, NAVY, bold=True)])
    ownership_bar(s, 0.85, y + 0.26, 11.6, 0.74, segs)
    y += 1.30
callout(s, 5.92, "Aim is not to minimise dilution — it is to maximise the "
        "value of the slice retained. A smaller share of a de-risked, larger "
        "company can be worth far more.", h=0.92)
notes(s, "The option pool is usually created pre-money — it dilutes founders, "
         "not the incoming investor.")

# ---- pre/post money formula + worked example ----
formula_slide(
    "Section 6  ·  Stakeholders",
    "Pre-money, post-money and the dilution math",
    [
        "Post-money = Pre-money + New Investment",
        "Investor % = New Investment / Post-money       "
        "Founder % after = before × ( Pre / Post )",
    ],
    [
        ("Pre-money", "venture value immediately before the new investment"),
        ("Post-money", "value immediately after — the basis for ownership %"),
        ("Note", "a pre-money option pool dilutes founders further"),
    ],
    [
        ("Raise US$2.5M at a US$4.4M pre-money", True),
        ("1 · Post-money = 4.4 + 2.5 = US$6.9M", False),
        ("2 · Investor % = 2.5 ÷ 6.9 = 36.2%", True),
        ("3 · Founders before = 100%", False),
        ("    → after = 100% × (4.4 ÷ 6.9) = 63.8%", True),
        ("4 · With a 10% pre-money option pool,", False),
        ("    founders fall further to ≈ 53.8%", False),
    ],
    "The same headline number implies different ownership depending on whether "
    "it is pre- or post-money — always confirm which.",
    notes_text="The pre- vs post-money confusion is a classic, costly error.")

# ======================================================================
# SECTION 7
# ======================================================================
divider("07", "Section Seven", "The Master Model",
        "Assembling the six methods into one integrated model — and "
        "comparing it with the other financial approaches.",
        "Section 7 is the synthesis.")

# ---- master model layer stack ----
s = diagram_slide("Section 7  ·  The master model",
                  "The master model — six methods, one argument")
layers = [
    ("LAYER 1", "Milestone roadmap", "defines the value-inflection timeline", STEP_COLORS[0]),
    ("LAYER 2", "Use of funds / rNPV", "bottom-up intrinsic value of next milestones", STEP_COLORS[1]),
    ("LAYER 3", "Benchmarking", "top-down, market-anchored exit value", STEP_COLORS[2]),
    ("LAYER 4", "Matching", "reconciles Layers 2 & 3 into a range", STEP_COLORS[3]),
    ("LAYER 5", "Uncertainty", "Golden-Triangle scenarios → a weighted range", STEP_COLORS[4]),
    ("LAYER 6", "Stakeholder terms", "enterprise value → value per stakeholder", STEP_COLORS[5]),
]
y = 1.34
for i, (tag, name, role, col) in enumerate(layers):
    indent = 0.85 + i * 0.30
    w = 11.6 - i * 0.60
    shp(s, RRECT, indent, y, w, 0.56, col, rad=0.16)
    tbox(s, indent + 0.20, y, 1.55, 0.56,
         [T(tag, 11.5, WHITE, bold=True)], anchor=MID)
    tbox(s, indent + 1.75, y, w - 1.95, 0.56,
         [T(name, 14, WHITE, bold=True, name=HEAD)], anchor=MID)
    tbox(s, indent + 1.75, y, w - 2.0, 0.56,
         [T(role, 11.5, RGBColor(0xEC, 0xF1, 0xF5), align=R)], anchor=MID)
    y += 0.605
shp(s, DARR, 6.56, y - 0.02, 0.24, 0.24, ACCENT)
ob = shp(s, RRECT, 2.4, y + 0.26, 8.5, 0.58, NAVY, rad=0.5)
shp_text(ob, [T("OUTPUT:  a valuation range  ·  a risk profile  ·  an "
                "ownership outcome", 13, WHITE, bold=True, align=C)])
callout(s, 6.10, "The master model is not a seventh method — it is the "
        "disciplined assembly of the previous six.", label="THE SYNTHESIS",
        h=0.62)
notes(s, "Each layer is one prior section.")

# ---- master model walkthrough ----
table_slide(
    "Section 7  ·  The master model",
    "Master model walkthrough — one venture, end to end",
    "Illustrative: a diagnostic-imaging venture carried through all six layers.",
    ["Layer", "Method applied", "Result"],
    [
        ["1  Roadmap", "TRL 4 → 9; four milestones to exit", "exit horizon: Year 5"],
        ["2  Use of funds", "US$2.5M funds M1+M2; CPOS 16.8%; r 25%", "intrinsic pre-money ≈ US$4.4M"],
        ["3  Benchmarking", "Liver-diagnostics TAM US$30.6B; comparables", "exit value ≈ US$80M"],
        ["4  Matching", "Bottom-up vs top-down reconciled", "defensible pre-money US$4–6M"],
        ["5  Uncertainty", "Down/Base/Up = 8 / 25 / 70; E[V] ≈ US$29M", "a range, not a point"],
        ["6  Stakeholder terms", "US$2.5M at US$4.5M pre; 1× pref", "investors ≈ 36%; founders keep majority"],
    ],
    [2.5, 5.3, 4.4], fsize=12.5,
    callout_text="Six layers, one argument: a valuation range, its risk "
                 "profile and the cap-table outcome — all traceable to "
                 "stated assumptions.",
    notes_text="Every figure traces back to a named assumption.")

# ---- comparison ----
table_slide(
    "Section 7  ·  Comparison",
    "Comparison — other financial valuation approaches",
    "The master model is the spine; these methods are cross-checks for different stages.",
    ["Approach", "Essence", "Best used when", "Caveat for deep tech"],
    [
        ["Revenue multiple", "EV = Revenue × multiple", "Early commercial revenue exists", "Silent pre-revenue; volatile"],
        ["Earnings / P/E", "Equity = Net profit × P/E", "Profitable, listed peers exist", "Not applicable while pre-profit"],
        ["EV/EBITDA & market multiples", "EV = EBITDA × peer multiple", "Scaling, near-peer financials", "Needs illiquidity / stage discounts"],
        ["DCF / rNPV", "Σ risk-adjusted CF ÷ (1+r)ₜ", "Any stage — the intrinsic anchor", "Highly assumption-sensitive"],
        ["VC method", "EV = Exit value ÷ target return", "Pricing a priced round", "Depends on exit assumptions"],
        ["Comparable transactions", "Apply observed deal metrics", "Reference deals exist", "Deal terms rarely disclosed"],
    ],
    [2.5, 3.0, 3.3, 3.6], fsize=12,
    callout_text="Cross-checks, not rivals. WBI was valued by DCF (WACC 12%, "
                 "terminal growth 3%) at US$55.6M, then sense-checked against "
                 "an industry P/E of 50–60×.",
    ctop=6.18,
    notes_text="The WBI case is a real DTL-adjacent valuation; capital-market "
               "multiples leverage listed peers but need discounts.")

# ---- method by stage ----
table_slide(
    "Section 7  ·  Comparison",
    "Choosing the right method by venture stage",
    "Method selection follows TRL stage — the matrix introduced in Section 1.",
    ["Stage / TRL", "Primary method", "Secondary cross-check"],
    [
        ["Concept (TRL 1–3)", "Milestone roadmap + qualitative scorecard", "Comparable seed rounds"],
        ["Prototype (TRL 4–6)", "Use-of-funds rNPV + VC method", "Benchmarked exit value"],
        ["Validation (TRL 7–8)", "rNPV + comparable transactions", "Revenue multiple, if pilot revenue"],
        ["Early commercial (TRL 9)", "Revenue & EV/EBITDA multiples", "DCF"],
        ["Scaling", "DCF + capital-market multiples", "P/E versus listed peers"],
    ],
    [3.0, 4.7, 4.5], fsize=13,
    callout_text="No single method is correct — the master model is correct "
                 "because it uses the right method at each stage.",
    notes_text="Stage drives method.")

# ======================================================================
# CLOSING
# ======================================================================
content_slide(
    "Closing  ·  Caveats", "Common pitfalls and academic caveats",
    [
        (0, "Single-point valuations — deep tech is a distribution; always "
            "present a range with its probabilities."),
        (0, "Confusing headline valuation with value received — terms and "
            "dilution can outweigh the price."),
        (0, "Double-counting risk — once a milestone is achieved, remove its "
            "POS from CPOS; do not also keep an inflated discount rate."),
        (0, "Optimistic POS — use base rates from comparable programmes, not "
            "the founder's confidence."),
        (0, "Stale benchmarks — venture references age quickly; date and "
            "adjust every comparable."),
        (0, "Ignoring the option pool — it is usually created pre-money and "
            "dilutes founders, not investors."),
    ],
    callout_text="A valuation is only as honest as its weakest assumption — "
                 "state assumptions explicitly so others can challenge them.",
    notes_text="These are the recurring errors in DTL venture reviews.")

# ---- workshop ----
s = diagram_slide("Closing  ·  Workshop",
                  "Workshop exercise — build a valuation for your venture")
tbox(s, 0.62, 1.26, 12.1, 0.4,
     [T("Work with your own venture, or a DTL case provided by the facilitator.",
        14, GREY, italic=True)])
steps = [
    ("1", "Draw the milestone roadmap from today's TRL to a credible exit."),
    ("2", "State the next round's use of funds and the milestones it buys."),
    ("3", "Find three market references for the exit value."),
    ("4", "Build the rNPV, then reconcile it with benchmarking (matching)."),
    ("5", "Write the downside and upside across the Golden Triangle."),
    ("6", "Set a pre-money, a raise amount, and show the resulting cap table."),
]
y = 1.84
for num, txt in steps:
    shp(s, OVAL, 0.85, y, 0.62, 0.62, STEP_COLORS[int(num) - 1])
    tbox(s, 0.85, y, 0.62, 0.62,
         [T(num, 18, WHITE, bold=True, name=HEAD, align=C)], anchor=MID)
    card = shp(s, RRECT, 1.70, y + 0.02, 10.6, 0.58, LIGHT, rad=0.18)
    shp_text(card, [T(txt, 14.5, DARK, name=HEAD)], ml=0.26)
    y += 0.72
callout(s, 6.22, "Deliverable: a one-page valuation memo — a value range, the "
        "six layers behind it, and every assumption stated.", h=0.74)
notes(s, "The one-page memo is the real test of whether the training landed.")

content_slide(
    "Closing  ·  Summary", "Key takeaways and further reading",
    [
        (0, "Valuation is an argument built from milestones, funds and market "
            "evidence — it is assembled, not guessed."),
        (0, "rNPV makes de-risking measurable; benchmarking anchors it to the "
            "market; matching reconciles the two."),
        (0, "Always express uncertainty as a downside/upside range across the "
            "Golden Triangle — technology, organisation, commercial."),
        (0, "Stakeholder terms decide who receives the value — read them as "
            "carefully as the headline number."),
        (1, "Damodaran, A. — Valuing Young, Start-up and Growth Companies (NYU Stern)."),
        (1, "Metrick & Yasuda — Venture Capital and the Finance of Innovation."),
        (1, "Razgaitis, R. — Valuation and Dealmaking of Technology-Based IP."),
        (1, "Deep-Tech Lab (DTL) venture files — internal references used here."),
    ],
    callout_text="All figures in this deck are illustrative worked examples for "
                 "training purposes — they are not investment advice.",
    notes_text="Restate the four takeaways; point to the reading list.")

# ======================================================================
# APPENDIX — Glossary of abbreviations
# ======================================================================
def glossary_slide(kicker, title, lead, entries):
    s = new_slide()
    header(s, kicker, title)
    tbox(s, 0.62, 1.24, 12.1, 0.5, [T(lead, 14, GREY, italic=True)])
    half = (len(entries) + 1) // 2
    cols = [entries[:half], entries[half:]]
    for ci, col in enumerate(cols):
        x = 0.66 + ci * 6.18
        shp(s, RECT, x, 1.88, 0.09, 4.55, ACCENT)
        tb = s.shapes.add_textbox(Inches(x + 0.26), Inches(1.80),
                                  Inches(5.62), Inches(4.95))
        tf = tb.text_frame
        tf.word_wrap = True
        for ei, (abbr, full, gloss) in enumerate(col):
            p = tf.paragraphs[0] if ei == 0 else tf.add_paragraph()
            p.space_before = Pt(0 if ei == 0 else 11)
            r = p.add_run(); r.text = abbr
            r.font.size = Pt(15); r.font.bold = True
            r.font.color.rgb = ACCENT; r.font.name = MONO
            r = p.add_run(); r.text = "   " + full
            r.font.size = Pt(14.5); r.font.bold = True
            r.font.color.rgb = NAVY; r.font.name = HEAD
            g = tf.add_paragraph()
            g.space_before = Pt(1)
            r = g.add_run(); r.text = gloss
            r.font.size = Pt(13); r.font.color.rgb = GREY; r.font.name = BODY
    footer(s)
    notes(s, "Glossary appendix — every abbreviation used in the deck spelled "
             "out in full. Refer participants here whenever a term is unclear.")
    return s


glossary_slide(
    "Appendix  ·  Glossary", "Abbreviations in full — valuation & finance",
    "Every abbreviation used in this deck, spelled out and briefly explained.",
    [
        ("TRL", "Technology Readiness Level",
         "Maturity scale from 1 (basic idea) to 9 (proven in use)."),
        ("POS", "Probability of Success",
         "The chance that a given milestone is achieved."),
        ("CPOS", "Cumulative Probability of Success",
         "Product of every milestone POS up to a given date."),
        ("NPV", "Net Present Value",
         "Future cash flows discounted to their value today."),
        ("rNPV", "risk-adjusted Net Present Value",
         "NPV weighted by CPOS; also called eNPV (expected NPV)."),
        ("DCF", "Discounted Cash Flow",
         "Valuation by discounting projected future cash flows."),
        ("CF", "Cash Flow",
         "CF in year t — the cash, or exit value, expected then."),
        ("r", "Discount rate",
         "Venture cost of capital — typically 15–40% early-stage."),
        ("WACC", "Weighted Average Cost of Capital",
         "Blended cost of equity and debt used to discount cash flows."),
        ("EV", "Enterprise Value",
         "The value of the whole business — debt plus equity."),
        ("EBITDA", "Earnings Before Interest, Tax, Depreciation & Amortisation",
         "A common proxy for operating cash earnings."),
        ("P/E", "Price-to-Earnings ratio",
         "Share price divided by earnings per share."),
        ("VC", "Venture Capital",
         "The ‘VC method’ values a venture back from its target exit."),
        ("E[V]", "Expected Value",
         "Probability-weighted average across the scenarios."),
    ])

glossary_slide(
    "Appendix  ·  Glossary", "Abbreviations in full — market, science & deal",
    "Continued — market-sizing, scientific, regulatory and programme terms.",
    [
        ("DTL", "Deep-Tech Lab",
         "The venture-development programme this training supports."),
        ("IP", "Intellectual Property",
         "Patents, know-how and freedom-to-operate position."),
        ("R&D", "Research and Development",
         "The core scientific and engineering work of the venture."),
        ("TAM", "Total Addressable Market",
         "Total demand if the product reached every possible user."),
        ("SAM", "Serviceable Available Market",
         "The slice the product and channel can realistically serve."),
        ("SOM", "Serviceable Obtainable Market",
         "The market share the venture can actually win near-term."),
        ("CAGR", "Compound Annual Growth Rate",
         "Smoothed year-on-year growth rate over a period."),
        ("GTM", "Go-To-Market",
         "Commercial launch — first sales and channel activity."),
        ("FDA", "US Food and Drug Administration",
         "The United States medical-products regulator."),
        ("510(k)", "FDA Premarket Notification",
         "Clearance route for many Class II medical devices."),
        ("MRI", "Magnetic Resonance Imaging",
         "Imaging used by the diagnostic-venture worked example."),
        ("WBI", "Wellbiz International",
         "Real venture used in the valuation-comparison example."),
        ("US$M", "US dollars, millions",
         "Currency unit used throughout the worked examples."),
    ])

# ----------------------------------------------------------------------
out = "Valuing_Deep_Tech_Ventures_DTL_Training.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
